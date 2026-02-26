#!/usr/bin/env python3
"""Build a LinkedIn carousel PDF and PPTX for Pop! miniature popcorn machine.

Format: 1080x1350 px (4:5 portrait) — optimized for mobile feed.
Uses reportlab for PDF and python-pptx for PPTX generation.
"""

import os
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor, white
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Page size: 4:5 ratio --
PW = 190 * mm
PH = 237.5 * mm

# -- Colors --
DARK_BG = "#1A1A2E"
CARD_BG = "#22223A"
CARD_BG_ALT = "#1E1E34"
ACCENT_ORANGE = "#FF8C00"
ACCENT_GREEN = "#4EC978"
ACCENT_RED = "#FF4545"
ACCENT_BLUE = "#009BF5"
ACCENT_PURPLE = "#A855F7"
ACCENT_YELLOW = "#FFD700"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY = "#BBBBCC"
SOFT_WHITE = "#F0F0F5"

M = 10 * mm  # standard margin


def bg(c, color=DARK_BG):
    c.setFillColor(HexColor(color))
    c.rect(0, 0, PW, PH, fill=1, stroke=0)


def accent_strip(c, color=ACCENT_ORANGE):
    c.setFillColor(HexColor(color))
    c.rect(0, PH - 2 * mm, PW, 2 * mm, fill=1, stroke=0)


def card(c, x, y, w, h, color=CARD_BG):
    """Draw a card. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    c.roundRect(x, PH - y - h, w, h, 3 * mm, fill=1, stroke=0)


def card_flat(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def txt(c, x, y, text, size=14, color=WHITE_HEX, bold=False, align="left", max_w=None):
    """Draw text. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    ry = PH - y - size * 0.35
    if align == "center" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + (max_w - tw) / 2, ry, text)
    elif align == "right" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + max_w - tw, ry, text)
    else:
        c.drawString(x, ry, text)


def txt_wrap(c, x, y, text, size=11, color=WHITE_HEX, bold=False, line_h=None, max_w=None):
    """Draw wrapped text. Returns y after last line."""
    if line_h is None:
        line_h = size * 1.4
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFillColor(HexColor(color))
    c.setFont(font, size)
    if max_w is None:
        max_w = PW - 2 * M
    words = text.split()
    lines = []
    current = ""
    for w in words:
        test = current + (" " if current else "") + w
        if c.stringWidth(test, font, size) > max_w:
            if current:
                lines.append(current)
            current = w
        else:
            current = test
    if current:
        lines.append(current)
    for line in lines:
        c.drawString(x, PH - y - size * 0.35, line)
        y += line_h
    return y


def bar(c, x, y, w, h, color):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def circle_num(c, x, y, num, color):
    r = 4 * mm
    c.setFillColor(HexColor(color))
    c.circle(x + r, PH - y - r, r, fill=1, stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    tw = c.stringWidth(str(num), "Helvetica-Bold", 12)
    c.drawString(x + r - tw / 2, PH - y - r - 4, str(num))


def footer(c, page, total):
    c.setFillColor(HexColor(LIGHT_GRAY))
    c.setFont("Helvetica", 8)
    label = f"{page}/{total}"
    tw = c.stringWidth(label, "Helvetica", 8)
    c.drawString(PW - M - tw, 4 * mm, label)


TOTAL_PAGES = 8
output = os.path.join(_DIR, "Pop_Carousel.pdf")
c = canvas.Canvas(output, pagesize=(PW, PH))

# ================================================================
# PAGE 1: Title
# ================================================================
bg(c)
accent_strip(c, ACCENT_YELLOW)

txt(c, M, 14 * mm, "Pop!", size=42, color=WHITE_HEX, bold=True)

bar(c, M, 26 * mm, 35 * mm, 1 * mm, ACCENT_YELLOW)

txt_wrap(c, M, 32 * mm,
         "One kernel at a time. Physics does the sorting.",
         size=18, color=ACCENT_GREEN, max_w=PW - 2 * M)

txt_wrap(c, M, 48 * mm,
         "A miniature desktop popcorn machine that pops kernels individually "
         "on a slow conveyor belt through a hot zone. Popped kernels escape "
         "into a bowl. Duds ride to the end. USB-C powered. Mesmerizing.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Image
img_path = os.path.join(_DIR, "cross_section_illustration_pop_miniature_popcorn_machine.png")
if os.path.exists(img_path):
    img_top_y = 72 * mm
    img_max_w = PW - 2 * M
    img_max_h = 135 * mm
    c.drawImage(ImageReader(img_path),
                M, PH - img_top_y - img_max_h,
                width=img_max_w, height=img_max_h,
                preserveAspectRatio=True, anchor='sw', mask='auto')

# Bottom bar
card_flat(c, 0, PH - 8 * mm, PW, 8 * mm, CARD_BG)
txt(c, M, PH - 5 * mm, "Product Overview  |  Concept Stage  |  2026",
    size=9, color=LIGHT_GRAY)

footer(c, 1, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 2: The Problem
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "The Problem", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Microwave popcorn is invisible",
     "Push a button, wait 3 minutes, listen for the slowdown. "
     "Nothing to see. No ritual. No spectacle. The bag arrives "
     "pre-made with artificial butter flavor and regret.",
     ACCENT_RED),
    ("Stovetop is overkill",
     "Oil, pot, lid, constant shaking, burned kernels stuck to "
     "the bottom, cleanup. All for a snack. Nobody does this "
     "at their desk.",
     ACCENT_ORANGE),
    ("No desk-scale popcorn experience exists",
     "There is nothing that turns popcorn into a meditative, "
     "kernel-by-kernel show you watch while you work. "
     "Pop! fills the space between snack and spectacle.",
     ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(problems):
    y_top = 30 * mm + i * 55 * mm
    card(c, M, y_top, PW - 2 * M, 49 * mm, CARD_BG)
    bar(c, M, y_top, PW - 2 * M, 1.5 * mm, color)
    txt(c, M + 5 * mm, y_top + 8 * mm, title, size=16, color=color, bold=True)
    txt_wrap(c, M + 5 * mm, y_top + 20 * mm, desc,
             size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Target users
card_flat(c, M, 200 * mm, PW - 2 * M, 30 * mm, CARD_BG)
txt(c, M + 4 * mm, 204 * mm, "TARGET USERS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 213 * mm, "Desk workers  |  Dorm dwellers  |  Gadget lovers",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 222 * mm, "Gift recipients  |  Snack enthusiasts  |  Ages 14+",
    size=12, color=WHITE_HEX, bold=True)

footer(c, 2, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 3: How It Works
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "How It Works", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

steps = [
    ("Fill", "Pour kernels into the hopper",
     "Lift the clear lid, pour in two tablespoons of kernels. "
     "The hopper holds ~100. Plug in the USB-C cable.",
     ACCENT_BLUE),
    ("Start", "Press the button, wait 45 seconds",
     "LED ring glows amber as the PTC heater preheats to 200C. "
     "When the ring turns green, the conveyor belt starts.",
     ACCENT_GREEN),
    ("Watch", "Kernels pop one at a time",
     "Kernels slide down the gravity chute single-file onto the belt. "
     "The belt carries each through the hot zone. Pop! The kernel "
     "explodes, escapes the belt, tumbles into the bowl.",
     ACCENT_YELLOW),
    ("Eat", "Popcorn bowl fills, duds go to waste",
     "Popped kernels escape sideways into the clear bowl. "
     "Duds ride the belt to the end and drop into the dud bowl. "
     "No sensors, no sorting logic. Physics does it.",
     ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(steps):
    y_top = 28 * mm + i * 50 * mm
    card(c, M, y_top, PW - 2 * M, 44 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, color)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=18, color=color, bold=True)
    txt(c, M + 16 * mm, y_top + 16 * mm, subtitle, size=11, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 26 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

footer(c, 3, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 4: Architecture
# ================================================================
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 12 * mm, "Architecture", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_BLUE)

# Kernel pipeline
card(c, M, 28 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 32 * mm, "Kernel Pipeline", size=11, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 41 * mm, "Hopper  -->  Chute  -->  Belt  -->  Hot Zone  -->  Bowl / Dud",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 49 * mm, "gravity       single-file    slow grip     200C PTC       physics sorting",
    size=8, color=LIGHT_GRAY)

# Key insight
card(c, M, 60 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 64 * mm, "Key Insight: No Sensors for Sorting", size=11, color=ACCENT_GREEN, bold=True)
txt(c, M + 4 * mm, 72 * mm, "Popped = 25mm, escapes belt.  Unpopped = 6mm, stays gripped.  Physics wins.",
    size=10, color=WHITE_HEX)

# Subsystems
bar(c, M, 88 * mm, PW - 2 * M, 1 * mm, ACCENT_BLUE)
txt(c, M, 92 * mm, "SUBSYSTEMS", size=11, color=ACCENT_BLUE, bold=True)

subsystems = [
    ("Double Conveyor Belt", "PTFE-fiberglass mesh, grips kernels through hot zone, 5-8 mm/s", ACCENT_YELLOW),
    ("PTC Ceramic Heater (20W)", "Self-regulating, 200C, mounted under belt path", ACCENT_RED),
    ("Gravity Chute", "Angled slide, narrows to 9mm, single-files kernels passively", ACCENT_GREEN),
    ("MCU (ESP32-C3)", "PID temp control, motor PWM, LED effects, optional BLE", ACCENT_BLUE),
    ("USB-C PD (30W)", "20V @ 1.5A via STUSB4500, no wall brick", ACCENT_ORANGE),
    ("LED Ring (WS2812B x8)", "Amber = preheat, green = running, blue = done, red = fault", ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(subsystems):
    y_top = 98 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 4 * mm, name, size=11, color=WHITE_HEX, bold=True)
    txt(c, M + 4 * mm, y_top + 11 * mm, desc, size=8, color=LIGHT_GRAY)
    bar(c, M, y_top, 1.5 * mm, 14 * mm, color)

# Firmware simplicity
card(c, M, 202 * mm, PW - 2 * M, 28 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 206 * mm, "Dumb firmware, smart mechanism", size=13, color=ACCENT_PURPLE, bold=True)
txt_wrap(c, M + 4 * mm, 216 * mm,
         "The firmware runs PID temperature control and constant-speed motor drive. "
         "No pop detection, no sorting logic, no jam recovery. The conveyor belt "
         "and physics handle everything the firmware doesn't need to know about.",
         size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 8 * mm, line_h=12)

footer(c, 4, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 5: Physics-Based Sorting
# ================================================================
bg(c)
accent_strip(c, ACCENT_YELLOW)

txt(c, M, 12 * mm, "Physics-Based Sorting", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_YELLOW)

txt_wrap(c, M, 28 * mm,
         "How does the machine know which kernels popped?",
         size=13, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Unpopped
card(c, M, 42 * mm, PW - 2 * M, 50 * mm, CARD_BG)
bar(c, M, 42 * mm, PW - 2 * M, 1.5 * mm, ACCENT_RED)
txt(c, M + 5 * mm, 48 * mm, "Unpopped Kernel (Dud)", size=16, color=ACCENT_RED, bold=True)
txt(c, M + 5 * mm, 60 * mm, "Size: ~6mm  |  Hard, round", size=13, color=WHITE_HEX)
txt_wrap(c, M + 5 * mm, 72 * mm,
         "Fits snugly between the two belts. Belt tension grips it. "
         "Rides the full belt path. Drops into the dud bowl at the end.",
         size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=13)

# Popped
card(c, M, 98 * mm, PW - 2 * M, 55 * mm, CARD_BG)
bar(c, M, 98 * mm, PW - 2 * M, 1.5 * mm, ACCENT_GREEN)
txt(c, M + 5 * mm, 104 * mm, "Popped Kernel", size=16, color=ACCENT_GREEN, bold=True)
txt(c, M + 5 * mm, 116 * mm, "Size: ~25mm  |  Expanded 10-15x  |  Irregular", size=13, color=WHITE_HEX)
txt_wrap(c, M + 5 * mm, 128 * mm,
         "Way too big for the belt gap. The expansion force pushes "
         "the spring-loaded belts apart. The kernel breaks free and "
         "escapes sideways into the popcorn bowl. No sensor needed.",
         size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=13)

# The mechanism
card(c, M, 160 * mm, PW - 2 * M, 68 * mm, CARD_BG_ALT)
txt(c, M + 5 * mm, 165 * mm, "The conveyor belt IS the sensor", size=13, color=ACCENT_YELLOW, bold=True)

details = [
    ("Belt gap: ~4mm at rest",
     "Spring-loaded rollers apply light tension. "
     "6mm kernels get gripped. 25mm popcorn can't be held."),
    ("Belt speed: 5-8 mm/s",
     "Each kernel spends 10-16 seconds in the hot zone. "
     "Fast enough to be interesting, slow enough to pop."),
    ("Belt material: PTFE-coated fiberglass",
     "Survives 260C continuous. Food-safe (FDA approved). "
     "Low-stick so popcorn doesn't adhere."),
]

ry = 178 * mm
for title, desc in details:
    txt(c, M + 5 * mm, ry, f"- {title}:", size=10, color=WHITE_HEX, bold=True)
    ry = txt_wrap(c, M + 8 * mm, ry + 12, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 16 * mm, line_h=12)
    ry += 3 * mm

footer(c, 5, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 6: Constraints & BOM
# ================================================================
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 12 * mm, "Constraints & BOM", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_ORANGE)

constraints = [
    ("Hot zone temp", "~200C", "PTC self-regulating, PID-controlled"),
    ("Belt speed", "5-8 mm/s", "10-16s per kernel in hot zone"),
    ("Power", "USB-C PD 30W", "19W steady state, 11W headroom"),
    ("Footprint", "150 x 180 x 140mm", "Fits on a desk next to monitor"),
    ("Weight", "~400g", "Light enough to move with one hand"),
    ("Food safety", "FDA / EU 1935/2004", "PTFE belt, Tritan bowls"),
]

for i, (name, value, note) in enumerate(constraints):
    y_top = 28 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=11, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, note, size=8, color=LIGHT_GRAY)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, value, size=11, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 8 * mm)

# BOM
bar(c, M, 133 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 137 * mm, "BOM ESTIMATE (1k units)", size=11, color=ACCENT_GREEN, bold=True)

bom = [
    ("ESP32-C3-MINI-1", "$1.50"),
    ("STUSB4500 PD controller", "$1.20"),
    ("USB-C + buck regulator", "$0.90"),
    ("MOSFET + DRV8837 motor driver", "$1.00"),
    ("PTC ceramic heater (20W)", "$1.50"),
    ("DC gearmotor (30 RPM)", "$2.00"),
    ("Conveyor belt assembly (belts + rollers + springs)", "$3.50"),
    ("Stainless steel chute", "$0.80"),
    ("NTC thermistor + passives", "$0.65"),
    ("WS2812B LEDs x8 + piezo", "$1.00"),
    ("PCB (50x40mm, 4-layer)", "$1.20"),
    ("Enclosure (hopper + housing + base + bowls)", "$5.10"),
    ("Thermal insulation", "$0.30"),
    ("Packaging + USB-C cable + sample kernels", "$1.50"),
    ("Assembly + test", "$2.50"),
]

for i, (item, cost) in enumerate(bom):
    y_top = 144 * mm + i * 6 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card_flat(c, M + 2 * mm, y_top, PW - 2 * M - 4 * mm, 5 * mm, bg_c)
    txt(c, M + 6 * mm, y_top + 1 * mm, item, size=7.5, color=SOFT_WHITE)
    txt(c, M + 6 * mm, y_top + 1 * mm, cost, size=7.5, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 16 * mm)

# Total
card_flat(c, M + 2 * mm, 234 * mm, PW - 2 * M - 4 * mm, 7 * mm, ACCENT_ORANGE)
txt(c, M + 6 * mm, 235.5 * mm, "Total COGS", size=9, color=WHITE_HEX, bold=True)
txt(c, M + 6 * mm, 235.5 * mm, "~$24.75  (5k: ~$18.50)", size=9, color=WHITE_HEX, bold=True,
    align="right", max_w=PW - 2 * M - 16 * mm)

footer(c, 6, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 7: Hardest Problems
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "Hardest Problems", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

hard_problems = [
    ("Belt grip vs. pop escape",
     "The belt must grip unpopped kernels (~6mm) firmly enough to "
     "transport them, yet release popped kernels (~25mm) reliably. "
     "Spring-loaded belt tension is the tuning variable. Too tight: "
     "popped kernels get crushed. Too loose: unpopped kernels slip. "
     "Must prototype with real kernels and real belts at M1."),
    ("Heat transfer through the belt",
     "Kernels are heated through PTFE-fiberglass mesh from a PTC "
     "heater below. If heat transfer is too slow, each kernel takes "
     ">12 seconds to pop and the cadence is boring. May need heater "
     "on both sides or direct radiant exposure through mesh openings. "
     "Target: pop within 8-12 seconds of entering the hot zone."),
    ("Gravity chute jamming",
     "The passive chute must single-file kernels from a bulk hopper "
     "using only geometry and gravity. Kernels are 4-8mm, irregular, "
     "and tend to bridge in narrowing channels. V-groove cross-section "
     "and polished surface help. A small vibration motor is the fallback. "
     "Must test with multiple kernel brands at M1."),
]

for i, (title, desc) in enumerate(hard_problems):
    y_top = 30 * mm + i * 62 * mm
    card(c, M, y_top, PW - 2 * M, 56 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, ACCENT_RED)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=13, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 18 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom
card_flat(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "All three are mechanism problems solved by prototyping, not simulation. "
         "M1 milestone: 3D-printed frame, off-shelf belts, PTC heater, real kernels.",
         size=11, color=ACCENT_ORANGE, max_w=PW - 2 * M - 8 * mm, line_h=14)

footer(c, 7, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 8: Gate Result & Open Items
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "Gate Result & Next", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

# Gate badge
card(c, M, 28 * mm, PW - 2 * M, 22 * mm, CARD_BG)
card(c, M + 4 * mm, 31 * mm, 50 * mm, 16 * mm, ACCENT_GREEN)
txt(c, M + 8 * mm, 35 * mm, "GATE: PASS", size=16, color=WHITE_HEX, bold=True)
txt(c, M + 8 * mm, 43 * mm, "70 pass / 13 N/A / 3 minor", size=9, color=WHITE_HEX)
txt(c, M + 60 * mm, 36 * mm, "System description complete.", size=11, color=WHITE_HEX)
txt(c, M + 60 * mm, 44 * mm, "Ready for mechanism prototype.", size=11, color=ACCENT_GREEN, bold=True)

# Power summary
card(c, M, 54 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 58 * mm, "POWER", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 66 * mm, "USB-C PD: 20V @ 1.5A (30W)  |  Steady state: ~19W",
    size=11, color=WHITE_HEX)
txt(c, M + 4 * mm, 73 * mm, "Heater 17W  |  Motor 0.3W  |  MCU 0.15W  |  LEDs 0.8W  |  11W headroom",
    size=10, color=LIGHT_GRAY)

# Key specs
card(c, M, 80 * mm, PW - 2 * M, 22 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 84 * mm, "KEY SPECS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 92 * mm, "150x180x140mm  |  ~400g  |  Clear enclosure  |  Tritan bowls",
    size=10, color=WHITE_HEX)
txt(c, M + 4 * mm, 99 * mm, "ESP32-C3  |  BLE 5.0  |  OTA  |  Retail $50-80  |  BOM <$25",
    size=10, color=LIGHT_GRAY)

# Open items
bar(c, M, 108 * mm, PW - 2 * M, 1 * mm, ACCENT_ORANGE)
txt(c, M, 112 * mm, "8 OPEN ITEMS", size=11, color=ACCENT_ORANGE, bold=True)

open_items = [
    ("M1", "Belt grip vs. pop escape tuning with real kernels"),
    ("M1", "Heat transfer rate through PTFE-fiberglass belt"),
    ("M1", "Gravity chute jamming across kernel brands"),
    ("M2", "Popped kernel escape direction / bowl placement"),
    ("M2", "Thermal isolation (outer surface <45C)"),
    ("M3", "Belt longevity at 200C (target 500+ hours)"),
    ("M2", "UL/ETL path for <25W heating appliance"),
    ("M2", "USB-C PD compatibility across charger brands"),
]

for i, (milestone, desc) in enumerate(open_items):
    y_top = 118 * mm + i * 12 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 10 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3 * mm, milestone, size=9, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 18 * mm, y_top + 3 * mm, desc, size=9, color=SOFT_WHITE)

# CTA
card(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "Next step: 3D-print the frame, source PTFE belts and a PTC heater, "
         "build the conveyor mechanism. Pop real kernels. Validate belt grip, "
         "chute flow, and escape geometry before committing to tooling.",
         size=10, color=WHITE_HEX, max_w=PW - 2 * M - 8 * mm, line_h=13)

footer(c, 8, TOTAL_PAGES)
c.showPage()

# ================================================================
c.save()
print(f"Saved {TOTAL_PAGES}-page carousel PDF to {output}")
print(f"Page size: {PW/mm:.0f} x {PH/mm:.0f} mm (4:5 ratio)")


# ================================================================
# PPTX GENERATION
# ================================================================
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(9.375)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# -- PPTX helper colors --
C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_CARD_BG = RGBColor(0x22, 0x22, 0x3A)
C_CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)
C_ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
C_ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
C_ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
C_ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_ACCENT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

PM = Inches(0.4)


def pptx_bg(slide, color=C_DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def pptx_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_text(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return txBox


def pptx_add_para(tf, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def pptx_footer(slide, page, total):
    pptx_text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.3),
              Inches(0.6), Inches(0.25), f"{page}/{total}",
              size=8, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ================================================================
# PPTX PAGE 1: Title
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_YELLOW)

pptx_text(slide, PM, Inches(0.5), SLIDE_W - 2 * PM, Inches(0.7),
          "Pop!", size=42, bold=True)

pptx_rect(slide, PM, Inches(1.1), Inches(1.4), Inches(0.04), C_ACCENT_YELLOW)

pptx_text(slide, PM, Inches(1.3), SLIDE_W - 2 * PM, Inches(0.4),
          "One kernel at a time. Physics does the sorting.", size=18, color=C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(1.85), SLIDE_W - 2 * PM, Inches(0.6),
          "A miniature desktop popcorn machine that pops kernels individually "
          "on a slow conveyor belt through a hot zone. Popped kernels escape "
          "into a bowl. Duds ride to the end. USB-C powered. Mesmerizing.",
          size=11, color=C_LIGHT_GRAY)

img_path = os.path.join(_DIR, "cross_section_illustration_pop_miniature_popcorn_machine.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, PM, Inches(2.8),
                             width=SLIDE_W - 2 * PM, height=Inches(5.3))

pptx_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_CARD_BG)
pptx_text(slide, PM, SLIDE_H - Inches(0.3), SLIDE_W - 2 * PM, Inches(0.2),
          "Product Overview  |  Concept Stage  |  2026", size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 1, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 2: The Problem
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "The Problem", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_problems = [
    ("Microwave popcorn is invisible",
     "Push a button, wait 3 minutes, listen for the slowdown. "
     "Nothing to see. No ritual. No spectacle. The bag arrives "
     "pre-made with artificial butter flavor and regret.",
     C_ACCENT_RED),
    ("Stovetop is overkill",
     "Oil, pot, lid, constant shaking, burned kernels stuck to "
     "the bottom, cleanup. All for a snack. Nobody does this "
     "at their desk.",
     C_ACCENT_ORANGE),
    ("No desk-scale popcorn experience exists",
     "There is nothing that turns popcorn into a meditative, "
     "kernel-by-kernel show you watch while you work. "
     "Pop! fills the space between snack and spectacle.",
     C_ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(pptx_problems):
    y = Inches(1.15 + i * 2.15)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
    pptx_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.06), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.3),
              title, size=16, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.65), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
              desc, size=11, color=C_LIGHT_GRAY)

pptx_rounded_rect(slide, PM, Inches(7.7), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
          "TARGET USERS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.1), SLIDE_W - 2 * PM, Inches(0.25),
          "Desk workers  |  Dorm dwellers  |  Gadget lovers", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Gift recipients  |  Snack enthusiasts  |  Ages 14+", size=12, bold=True)

pptx_footer(slide, 2, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 3: How It Works
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "How It Works", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

pptx_steps = [
    ("1. Fill", "Pour kernels into the hopper",
     "Lift the clear lid, pour in two tablespoons of kernels. "
     "The hopper holds ~100. Plug in the USB-C cable.",
     C_ACCENT_BLUE),
    ("2. Start", "Press the button, wait 45 seconds",
     "LED ring glows amber as the PTC heater preheats to 200C. "
     "When the ring turns green, the conveyor belt starts.",
     C_ACCENT_GREEN),
    ("3. Watch", "Kernels pop one at a time",
     "Kernels slide down the gravity chute single-file onto the belt. "
     "The belt carries each through the hot zone. Pop! The kernel "
     "explodes, escapes the belt, tumbles into the bowl.",
     C_ACCENT_YELLOW),
    ("4. Eat", "Popcorn bowl fills, duds go to waste",
     "Popped kernels escape sideways into the clear bowl. "
     "Duds ride the belt to the end and drop into the dud bowl. "
     "No sensors, no sorting logic. Physics does it.",
     C_ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(pptx_steps):
    y = Inches(1.1 + i * 1.95)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.35),
              title, size=18, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.45), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              subtitle, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.8), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_footer(slide, 3, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 4: Architecture
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Architecture", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_BLUE)

# Kernel pipeline
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(1.15), SLIDE_W - 2 * PM, Inches(0.2),
          "Kernel Pipeline", size=11, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Hopper  -->  Chute  -->  Belt  -->  Hot Zone  -->  Bowl / Dud",
          size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.8), SLIDE_W - 2 * PM, Inches(0.2),
          "gravity       single-file    slow grip     200C PTC       physics sorting",
          size=8, color=C_LIGHT_GRAY)

# Key insight
pptx_rounded_rect(slide, PM, Inches(2.35), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Key Insight: No Sensors for Sorting", size=11, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.7), SLIDE_W - 2 * PM, Inches(0.25),
          "Popped = 25mm, escapes belt.  Unpopped = 6mm, stays gripped.  Physics wins.",
          size=10, bold=True)

pptx_rect(slide, PM, Inches(3.4), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_BLUE)
pptx_text(slide, PM, Inches(3.5), SLIDE_W - 2 * PM, Inches(0.2),
          "SUBSYSTEMS", size=11, color=C_ACCENT_BLUE, bold=True)

pptx_subsystems = [
    ("Double Conveyor Belt", "PTFE-fiberglass mesh, grips kernels through hot zone, 5-8 mm/s", C_ACCENT_YELLOW),
    ("PTC Ceramic Heater (20W)", "Self-regulating, 200C, mounted under belt path", C_ACCENT_RED),
    ("Gravity Chute", "Angled slide, narrows to 9mm, single-files kernels passively", C_ACCENT_GREEN),
    ("MCU (ESP32-C3)", "PID temp control, motor PWM, LED effects, optional BLE", C_ACCENT_BLUE),
    ("USB-C PD (30W)", "20V @ 1.5A via STUSB4500, no wall brick", C_ACCENT_ORANGE),
    ("LED Ring (WS2812B x8)", "Amber = preheat, green = running, blue = done, red = fault", C_ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(pptx_subsystems):
    y = Inches(3.8 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_rect(slide, PM, y, Inches(0.06), Inches(0.55), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), SLIDE_W - 2 * PM, Inches(0.2),
              name, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.28), SLIDE_W - 2 * PM, Inches(0.2),
              desc, size=8, color=C_LIGHT_GRAY)

# Firmware simplicity
pptx_rounded_rect(slide, PM, Inches(7.75), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.25),
          "Dumb firmware, smart mechanism", size=13, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.15), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.6),
          "The firmware runs PID temperature control and constant-speed motor drive. "
          "No pop detection, no sorting logic, no jam recovery. The conveyor belt "
          "and physics handle everything the firmware doesn't need to know about.",
          size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 4, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 5: Physics-Based Sorting
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_YELLOW)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Physics-Based Sorting", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_YELLOW)

pptx_text(slide, PM, Inches(1.05), SLIDE_W - 2 * PM, Inches(0.4),
          "How does the machine know which kernels popped?",
          size=13, color=C_LIGHT_GRAY)

# Unpopped
pptx_rounded_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
pptx_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_RED)
pptx_text(slide, PM + Inches(0.2), Inches(1.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Unpopped Kernel (Dud)", size=16, color=C_ACCENT_RED, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(2.2), SLIDE_W - 2 * PM, Inches(0.25),
          "Size: ~6mm  |  Hard, round", size=13)
pptx_text(slide, PM + Inches(0.2), Inches(2.6), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.7),
          "Fits snugly between the two belts. Belt tension grips it. "
          "Rides the full belt path. Drops into the dud bowl at the end.",
          size=10, color=C_LIGHT_GRAY)

# Popped
pptx_rounded_rect(slide, PM, Inches(3.7), SLIDE_W - 2 * PM, Inches(2.1), C_CARD_BG)
pptx_rect(slide, PM, Inches(3.7), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.2), Inches(3.85), SLIDE_W - 2 * PM, Inches(0.3),
          "Popped Kernel", size=16, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(4.3), SLIDE_W - 2 * PM, Inches(0.25),
          "Size: ~25mm  |  Expanded 10-15x  |  Irregular", size=13)
pptx_text(slide, PM + Inches(0.2), Inches(4.7), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
          "Way too big for the belt gap. The expansion force pushes "
          "the spring-loaded belts apart. The kernel breaks free and "
          "escapes sideways into the popcorn bowl. No sensor needed.",
          size=10, color=C_LIGHT_GRAY)

# Mechanism details
pptx_rounded_rect(slide, PM, Inches(6.1), SLIDE_W - 2 * PM, Inches(2.75), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.2), Inches(6.2), SLIDE_W - 2 * PM, Inches(0.25),
          "The conveyor belt IS the sensor", size=13, color=C_ACCENT_YELLOW, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(6.6), SLIDE_W - 2 * PM - Inches(0.4), Inches(2.0),
                  "Belt gap: ~4mm at rest — spring-loaded rollers apply light tension. "
                  "6mm kernels get gripped. 25mm popcorn can't be held.", size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
pptx_add_para(tf, "Belt speed: 5-8 mm/s — each kernel spends 10-16 seconds in the hot zone. "
              "Fast enough to be interesting, slow enough to pop.",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Belt material: PTFE-coated fiberglass — survives 260C continuous. "
              "Food-safe (FDA approved). Low-stick so popcorn doesn't adhere.",
              size=9, color=C_LIGHT_GRAY, space_before=6)

pptx_footer(slide, 5, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 6: Constraints & BOM
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Constraints & BOM", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_ORANGE)

pptx_constraints = [
    ("Hot zone temp", "~200C", "PTC self-regulating, PID-controlled"),
    ("Belt speed", "5-8 mm/s", "10-16s per kernel in hot zone"),
    ("Power", "USB-C PD 30W", "19W steady state, 11W headroom"),
    ("Footprint", "150 x 180 x 140mm", "Fits on a desk next to monitor"),
    ("Weight", "~400g", "Light enough to move with one hand"),
    ("Food safety", "FDA / EU 1935/2004", "PTFE belt, Tritan bowls"),
]

for i, (name, value, note) in enumerate(pptx_constraints):
    y = Inches(1.1 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), Inches(3), Inches(0.2),
              name, size=11, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.3), Inches(5), Inches(0.2),
              note, size=8, color=C_LIGHT_GRAY)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.2),
              value, size=11, bold=True, align=PP_ALIGN.RIGHT)

# BOM
pptx_rect(slide, PM, Inches(5.15), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(5.25), SLIDE_W - 2 * PM, Inches(0.2),
          "BOM ESTIMATE (1k units)", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_bom = [
    ("ESP32-C3-MINI-1", "$1.50"),
    ("STUSB4500 PD controller", "$1.20"),
    ("USB-C + buck regulator", "$0.90"),
    ("MOSFET + DRV8837 motor driver", "$1.00"),
    ("PTC ceramic heater (20W)", "$1.50"),
    ("DC gearmotor (30 RPM)", "$2.00"),
    ("Conveyor belt assembly", "$3.50"),
    ("Stainless steel chute", "$0.80"),
    ("NTC thermistor + passives", "$0.65"),
    ("WS2812B LEDs x8 + piezo", "$1.00"),
    ("PCB (50x40mm, 4-layer)", "$1.20"),
    ("Enclosure (all parts)", "$5.10"),
    ("Thermal insulation", "$0.30"),
    ("Packaging + cable + samples", "$1.50"),
    ("Assembly + test", "$2.50"),
]

for i, (item, cost) in enumerate(pptx_bom):
    y = Inches(5.5 + i * 0.24)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rect(slide, PM + Inches(0.08), y, SLIDE_W - 2 * PM - Inches(0.16), Inches(0.2), bg_c)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.01), Inches(4), Inches(0.16),
              item, size=7.5, color=C_SOFT_WHITE)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.01), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.16),
              cost, size=7.5, bold=True, align=PP_ALIGN.RIGHT)

# Total
pptx_rect(slide, PM + Inches(0.08), Inches(9.12), SLIDE_W - 2 * PM - Inches(0.16), Inches(0.25), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.2), Inches(9.13), Inches(2), Inches(0.2),
          "Total COGS", size=9, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(9.13), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.2),
          "~$24.75  (5k: ~$18.50)", size=9, bold=True, align=PP_ALIGN.RIGHT)

pptx_footer(slide, 6, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 7: Hardest Problems
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Hardest Problems", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_hard = [
    ("Belt grip vs. pop escape",
     "The belt must grip unpopped kernels (~6mm) firmly enough to "
     "transport them, yet release popped kernels (~25mm) reliably. "
     "Spring-loaded belt tension is the tuning variable. Too tight: "
     "popped kernels get crushed. Too loose: unpopped kernels slip. "
     "Must prototype with real kernels and real belts at M1."),
    ("Heat transfer through the belt",
     "Kernels are heated through PTFE-fiberglass mesh from a PTC "
     "heater below. If heat transfer is too slow, each kernel takes "
     ">12 seconds to pop and the cadence is boring. May need heater "
     "on both sides or direct radiant exposure through mesh openings. "
     "Target: pop within 8-12 seconds of entering the hot zone."),
    ("Gravity chute jamming",
     "The passive chute must single-file kernels from a bulk hopper "
     "using only geometry and gravity. Kernels are 4-8mm, irregular, "
     "and tend to bridge in narrowing channels. V-groove cross-section "
     "and polished surface help. A small vibration motor is the fallback. "
     "Must test with multiple kernel brands at M1."),
]

for i, (title, desc) in enumerate(pptx_hard):
    y = Inches(1.15 + i * 2.4)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.15), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              f"{i+1}. {title}", size=13, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.55), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.4),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_rounded_rect(slide, PM, Inches(8.4), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "All three are mechanism problems solved by prototyping, not simulation. "
          "M1: 3D-printed frame, off-shelf belts, PTC heater, real kernels.",
          size=11, color=C_ACCENT_ORANGE)

pptx_footer(slide, 7, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 8: Gate Result & Next
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Gate Result & Next", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

# Gate badge
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_rect(slide, PM + Inches(0.15), Inches(1.2), Inches(2.0), Inches(0.6), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.3), Inches(1.22), Inches(1.8), Inches(0.3),
          "GATE: PASS", size=16, bold=True)
pptx_text(slide, PM + Inches(0.3), Inches(1.52), Inches(1.8), Inches(0.2),
          "70 pass / 13 N/A / 3 minor", size=9)
pptx_text(slide, PM + Inches(2.4), Inches(1.28), Inches(4), Inches(0.2),
          "System description complete.", size=11)
pptx_text(slide, PM + Inches(2.4), Inches(1.55), Inches(4), Inches(0.2),
          "Ready for mechanism prototype.", size=11, color=C_ACCENT_GREEN, bold=True)

# Power summary
pptx_rounded_rect(slide, PM, Inches(2.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.15), SLIDE_W - 2 * PM, Inches(0.2),
          "POWER", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "USB-C PD: 20V @ 1.5A (30W)  |  Steady state: ~19W", size=11)
pptx_text(slide, PM + Inches(0.15), Inches(2.65), SLIDE_W - 2 * PM, Inches(0.2),
          "Heater 17W  |  Motor 0.3W  |  MCU 0.15W  |  LEDs 0.8W  |  11W headroom",
          size=10, color=C_LIGHT_GRAY)

# Key specs
pptx_rounded_rect(slide, PM, Inches(3.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(3.15), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY SPECS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(3.4), SLIDE_W - 2 * PM, Inches(0.2),
          "150x180x140mm  |  ~400g  |  Clear enclosure  |  Tritan bowls", size=10)
pptx_text(slide, PM + Inches(0.15), Inches(3.65), SLIDE_W - 2 * PM, Inches(0.2),
          "ESP32-C3  |  BLE 5.0  |  OTA  |  Retail $50-80  |  BOM <$25",
          size=10, color=C_LIGHT_GRAY)

# Open items
pptx_rect(slide, PM, Inches(4.2), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_ORANGE)
pptx_text(slide, PM, Inches(4.3), SLIDE_W - 2 * PM, Inches(0.2),
          "8 OPEN ITEMS", size=11, color=C_ACCENT_ORANGE, bold=True)

pptx_open_items = [
    ("M1", "Belt grip vs. pop escape tuning with real kernels"),
    ("M1", "Heat transfer rate through PTFE-fiberglass belt"),
    ("M1", "Gravity chute jamming across kernel brands"),
    ("M2", "Popped kernel escape direction / bowl placement"),
    ("M2", "Thermal isolation (outer surface <45C)"),
    ("M3", "Belt longevity at 200C (target 500+ hours)"),
    ("M2", "UL/ETL path for <25W heating appliance"),
    ("M2", "USB-C PD compatibility across charger brands"),
]

for i, (milestone, desc) in enumerate(pptx_open_items):
    y = Inches(4.6 + i * 0.47)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.38), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.07), Inches(0.5), Inches(0.2),
              milestone, size=9, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.7), y + Inches(0.07), SLIDE_W - 2 * PM - Inches(1.0), Inches(0.2),
              desc, size=9, color=C_SOFT_WHITE)

# CTA
pptx_rounded_rect(slide, PM, Inches(8.45), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.5), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "Next step: 3D-print the frame, source PTFE belts and a PTC heater, "
          "build the conveyor mechanism. Pop real kernels. Validate belt grip, "
          "chute flow, and escape geometry before committing to tooling.",
          size=10)

pptx_footer(slide, 8, TOTAL_PAGES)

# Save PPTX
pptx_output = os.path.join(_DIR, "Pop_Carousel.pptx")
prs.save(pptx_output)
print(f"Saved {TOTAL_PAGES}-slide carousel PPTX to {pptx_output}")
