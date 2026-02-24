#!/usr/bin/env python3
"""Build a LinkedIn carousel PDF and PPTX for Consumable Electric Toothbrush.

Format: 1080x1350 px (4:5 portrait) — optimized for mobile feed.
Uses reportlab for PDF and python-pptx for PPTX generation.
"""

import os
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor, white
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Page size: 4:5 ratio --
PW = 190 * mm
PH = 237.5 * mm

# -- Colors --
DARK_BG = "#1A1A2E"
CARD_BG = "#22223A"
CARD_BG_ALT = "#1E1E34"
ACCENT_ORANGE = "#FF8C00"
ACCENT_GREEN = "#4EC978"
ACCENT_RED = "#FF4545"
ACCENT_BLUE = "#009BF5"
ACCENT_PURPLE = "#A855F7"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY = "#BBBBCC"
SOFT_WHITE = "#F0F0F5"

M = 10 * mm  # standard margin


# ── PDF helpers ──────────────────────────────────────────────

def bg(c, color=DARK_BG):
    c.setFillColor(HexColor(color))
    c.rect(0, 0, PW, PH, fill=1, stroke=0)


def accent_strip(c, color=ACCENT_ORANGE):
    c.setFillColor(HexColor(color))
    c.rect(0, PH - 2 * mm, PW, 2 * mm, fill=1, stroke=0)


def card(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.roundRect(x, PH - y - h, w, h, 3 * mm, fill=1, stroke=0)


def card_flat(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def txt(c, x, y, text, size=14, color=WHITE_HEX, bold=False, align="left", max_w=None):
    c.setFillColor(HexColor(color))
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    ry = PH - y - size * 0.35
    if align == "center" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + (max_w - tw) / 2, ry, text)
    elif align == "right" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + max_w - tw, ry, text)
    else:
        c.drawString(x, ry, text)


def txt_wrap(c, x, y, text, size=11, color=WHITE_HEX, bold=False, line_h=None, max_w=None):
    if line_h is None:
        line_h = size * 1.4
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFillColor(HexColor(color))
    c.setFont(font, size)
    if max_w is None:
        max_w = PW - 2 * M
    words = text.split()
    lines, current = [], ""
    for w in words:
        test = current + (" " if current else "") + w
        if c.stringWidth(test, font, size) > max_w:
            if current:
                lines.append(current)
            current = w
        else:
            current = test
    if current:
        lines.append(current)
    for line in lines:
        c.drawString(x, PH - y - size * 0.35, line)
        y += line_h
    return y


def bar(c, x, y, w, h, color):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def circle_num(c, x, y, num, color):
    r = 4 * mm
    c.setFillColor(HexColor(color))
    c.circle(x + r, PH - y - r, r, fill=1, stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    tw = c.stringWidth(str(num), "Helvetica-Bold", 12)
    c.drawString(x + r - tw / 2, PH - y - r - 4, str(num))


def footer(c, page, total):
    c.setFillColor(HexColor(LIGHT_GRAY))
    c.setFont("Helvetica", 8)
    label = f"{page}/{total}"
    tw = c.stringWidth(label, "Helvetica", 8)
    c.drawString(PW - M - tw, 4 * mm, label)


# ═══════════════════════════════════════════════════════════════
# PDF GENERATION
# ═══════════════════════════════════════════════════════════════
TOTAL_PAGES = 8
output = os.path.join(_DIR, "Consumable_Toothbrush_Carousel.pdf")
c = canvas.Canvas(output, pagesize=(PW, PH))

# ── PAGE 1: Title ────────────────────────────────────────────
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 14 * mm, "Consumable Electric", size=34, color=WHITE_HEX, bold=True)
txt(c, M, 28 * mm, "Toothbrush", size=34, color=WHITE_HEX, bold=True)

bar(c, M, 37 * mm, 35 * mm, 1 * mm, ACCENT_BLUE)

txt_wrap(c, M, 43 * mm,
         "Use it. Toss it. Open a new one.",
         size=18, color=ACCENT_GREEN, max_w=PW - 2 * M)

txt_wrap(c, M, 56 * mm,
         "A battery-powered vibrating toothbrush designed as a true "
         "consumable. AAA alkaline, no MCU, dual-injection sealed. "
         "$3-5 retail vs. Oral-B Pulsar at $8-12.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M)

img_path = os.path.join(_DIR, "cross_section_illustration_consumable_toothbrush.png")
if os.path.exists(img_path):
    c.drawImage(ImageReader(img_path),
                M, PH - 75 * mm - 135 * mm,
                width=PW - 2 * M, height=135 * mm,
                preserveAspectRatio=True, anchor='sw', mask='auto')

card_flat(c, 0, PH - 8 * mm, PW, 8 * mm, CARD_BG)
txt(c, M, PH - 5 * mm, "Product Overview  |  Concept Stage  |  2026",
    size=9, color=LIGHT_GRAY)

footer(c, 1, TOTAL_PAGES)
c.showPage()

# ── PAGE 2: The Problem ──────────────────────────────────────
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "The Problem", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Oral-B Pulsar is overpriced for a disposable",
     "At $8-12, the Pulsar pretends to be reusable but has a sealed, "
     "non-replaceable battery. When it dies, you throw it away anyway. "
     "You're paying premium prices for a disposable product.",
     ACCENT_RED),
    ("Manual brushes don't clean as well",
     "Powered vibration removes significantly more plaque than manual "
     "brushing. But rechargeable electric brushes cost $30-200 plus "
     "replacement heads. Budget consumers are stuck with manual.",
     ACCENT_ORANGE),
    ("No honest sub-$5 powered option exists",
     "The market has premium rechargeable brushes ($30+) and the "
     "Oral-B Pulsar ($8-12). Nothing fills the gap below $5 -- "
     "a price point where powered brushing becomes truly disposable.",
     ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    y_top = 30 * mm + i * 55 * mm
    card(c, M, y_top, PW - 2 * M, 49 * mm, CARD_BG)
    bar(c, M, y_top, PW - 2 * M, 1.5 * mm, color)
    txt(c, M + 5 * mm, y_top + 8 * mm, title, size=14, color=color, bold=True)
    txt_wrap(c, M + 5 * mm, y_top + 20 * mm, desc,
             size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

card_flat(c, M, 200 * mm, PW - 2 * M, 30 * mm, CARD_BG)
txt(c, M + 4 * mm, 204 * mm, "TARGET USERS", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 213 * mm, "Budget-conscious consumers  |  Travelers",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 222 * mm, "Hotels & hospitality  |  Multi-pack buyers",
    size=12, color=WHITE_HEX, bold=True)

footer(c, 2, TOTAL_PAGES)
c.showPage()

# ── PAGE 3: How It Works ─────────────────────────────────────
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "How It Works", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

steps = [
    ("Open", "Tear open blister pack",
     "AAA alkaline battery is pre-installed with years of shelf life. "
     "No charging, no pairing, no setup of any kind.",
     ACCENT_BLUE),
    ("Press", "Push the button once",
     "Latching push-button switch turns motor on. ERM vibration motor "
     "spins up instantly. Bristles oscillate via the split-head mechanism.",
     ACCENT_GREEN),
    ("Brush", "2 minutes, twice daily",
     "Split bristle head concentrates oscillation at the bristle tips "
     "for effective plaque removal. Rinse under tap after use. IPX5 sealed.",
     ACCENT_PURPLE),
    ("Toss", "~90 days later, motor stops",
     "Alkaline cell drops below motor stall voltage. Clear end-of-life: "
     "it works or it doesn't. Discard and open a new one.",
     ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(steps):
    y_top = 28 * mm + i * 50 * mm
    card(c, M, y_top, PW - 2 * M, 44 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, color)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=18, color=color, bold=True)
    txt(c, M + 16 * mm, y_top + 16 * mm, subtitle, size=11, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 26 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

footer(c, 3, TOTAL_PAGES)
c.showPage()

# ── PAGE 4: Architecture ─────────────────────────────────────
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 12 * mm, "Architecture", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_BLUE)

# Power chain
card(c, M, 28 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 32 * mm, "Power Path", size=11, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 41 * mm, "AAA Cell --> Latching Switch --> ERM Motor",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 49 * mm, "1.5V alkaline      push on/off       60-100mA, direct drive",
    size=8, color=LIGHT_GRAY)

# Mechanical chain
card(c, M, 60 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 64 * mm, "Mechanical Path", size=11, color=ACCENT_GREEN, bold=True)
txt(c, M + 4 * mm, 72 * mm, "Motor --> Eccentric --> Linkage --> Split Head --> Bristles",
    size=10, color=WHITE_HEX, bold=True)

# Subsystems
bar(c, M, 88 * mm, PW - 2 * M, 1 * mm, ACCENT_BLUE)
txt(c, M, 92 * mm, "SUBSYSTEMS", size=11, color=ACCENT_BLUE, bold=True)

subsystems = [
    ("AAA Alkaline Cell (1.5V)", "1000-1200 mAh, ~750 min runtime, years of shelf life", ACCENT_ORANGE),
    ("Latching Push-Button", "Mechanical push-on/push-off, sealed by TPE overmold", ACCENT_GREEN),
    ("Cylindrical ERM Motor", "6x12mm, 60-100mA at 1.5V, axial mount in handle neck", ACCENT_PURPLE),
    ("Split Bristle Head", "Fixed + moving halves, living hinge pivot, 1-2mm oscillation", ACCENT_GREEN),
    ("Rigid PP Body (1st shot)", "Handle, motor pocket, battery tube, insert-molded bristles", ACCENT_BLUE),
    ("TPE Overmold (2nd shot)", "Seals all penetrations, grip texture, button membrane, head boot", ACCENT_ORANGE),
    ("Battery Cap + O-ring", "Threaded PP cap, static radial seal, factory-installed", ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(subsystems):
    y_top = 98 * mm + i * 15 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 13 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=10, color=WHITE_HEX, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, desc, size=8, color=LIGHT_GRAY)
    bar(c, M, y_top, 1.5 * mm, 13 * mm, color)

# No software callout
card(c, M, 206 * mm, PW - 2 * M, 24 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 210 * mm, "Zero software. Zero electronics.", size=13, color=ACCENT_RED, bold=True)
txt_wrap(c, M + 4 * mm, 220 * mm,
         "No MCU, no PCB, no firmware, no app, no cloud. The entire "
         "electrical system is 3 components and 2 wires.",
         size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 8 * mm, line_h=12)

footer(c, 4, TOTAL_PAGES)
c.showPage()

# ── PAGE 5: Key Innovation — Dual-Injection Sealing ──────────
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 12 * mm, "Dual-Injection Sealing", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 28 * mm,
         "How do you waterproof a $1.50 product for 90 days of wet use?",
         size=13, color=LIGHT_GRAY, max_w=PW - 2 * M)

# 1st shot
card(c, M, 42 * mm, PW - 2 * M, 45 * mm, CARD_BG)
bar(c, M, 42 * mm, PW - 2 * M, 1.5 * mm, ACCENT_BLUE)
txt(c, M + 5 * mm, 48 * mm, "1st Shot: Rigid PP Body", size=16, color=ACCENT_BLUE, bold=True)
txt_wrap(c, M + 5 * mm, 60 * mm,
         "Injection-mold the structural handle with bristles pre-loaded "
         "in the mold cavity. Plastic flows around each bristle base, "
         "creating a watertight seal at every tuft. No secondary sealing.",
         size=11, color=WHITE_HEX, max_w=PW - 2 * M - 10 * mm, line_h=14)

# 2nd shot
card(c, M, 94 * mm, PW - 2 * M, 50 * mm, CARD_BG)
bar(c, M, 94 * mm, PW - 2 * M, 1.5 * mm, ACCENT_GREEN)
txt(c, M + 5 * mm, 100 * mm, "2nd Shot: TPE Overmold", size=16, color=ACCENT_GREEN, bold=True)
txt_wrap(c, M + 5 * mm, 112 * mm,
         "After motor, wiring, switch, and battery are assembled into "
         "the rigid body, the 2nd injection encapsulates everything. "
         "TPE bonds chemically to PP, sealing all penetrations in one "
         "step: grip texture, button membrane, and dynamic head boot.",
         size=11, color=WHITE_HEX, max_w=PW - 2 * M - 10 * mm, line_h=14)

# 4 seal zones
card(c, M, 152 * mm, PW - 2 * M, 75 * mm, CARD_BG_ALT)
txt(c, M + 5 * mm, 157 * mm, "4 Seal Zones", size=13, color=ACCENT_ORANGE, bold=True)

seals = [
    ("1. Bristle insert-mold", "PP flows around nylon tuft bases during 1st injection"),
    ("2. TPE body overmold", "2nd shot seals handle, switch, wire penetrations"),
    ("3. TPE head boot", "Flexible boot around split-head pivot, survives millions of flex cycles"),
    ("4. O-ring battery cap", "Threaded PP cap compresses O-ring at handle base"),
]

ry = 169 * mm
for title, desc in seals:
    txt(c, M + 5 * mm, ry, title, size=10, color=WHITE_HEX, bold=True)
    ry = txt_wrap(c, M + 8 * mm, ry + 12, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 16 * mm, line_h=12)
    ry += 3 * mm

footer(c, 5, TOTAL_PAGES)
c.showPage()

# ── PAGE 6: Constraints & BOM ────────────────────────────────
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 12 * mm, "Constraints & BOM", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_ORANGE)

constraints = [
    ("BOM cost", "< $1.50", "At 10k units; < $1.20 at 50k"),
    ("Motor runtime", ">= 350 min", "~750 min calculated (AAA @ 80mA)"),
    ("Water resistance", "IPX5", "Dual-injection + O-ring, 90 days wet use"),
    ("Handle diameter", "<= 16mm", "AAA cell (10.5mm) + PP wall + TPE"),
    ("Retail price", "$3 - $5", "Undercutting Oral-B Pulsar ($8-12)"),
    ("Shelf life", "> 2 years", "Alkaline chemistry, sealed blister pack"),
]

for i, (name, value, note) in enumerate(constraints):
    y_top = 28 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=11, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, note, size=8, color=LIGHT_GRAY)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, value, size=11, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 8 * mm)

# BOM
bar(c, M, 133 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 137 * mm, "BOM ESTIMATE (10k units)", size=11, color=ACCENT_GREEN, bold=True)

bom = [
    ("ERM motor (cylindrical, 6x12mm)", "$0.12"),
    ("AAA alkaline cell", "$0.07"),
    ("Latching push-button switch", "$0.03"),
    ("Wiring (2x, 26AWG, tinned)", "$0.01"),
    ("Spring contact + plate contact", "$0.02"),
    ("PA-612 nylon bristle tufts (x35)", "$0.02"),
    ("O-ring (battery cap seal)", "$0.01"),
    ("PP resin (1st shot, ~8g)", "$0.02"),
    ("TPE resin (2nd shot, ~3g)", "$0.02"),
    ("Battery cap (threaded PP)", "$0.02"),
    ("Blister packaging + card", "$0.08"),
    ("Assembly + test (90-100s/unit)", "$0.72"),
]

for i, (item, cost) in enumerate(bom):
    y_top = 144 * mm + i * 7 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card_flat(c, M + 2 * mm, y_top, PW - 2 * M - 4 * mm, 6 * mm, bg_c)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, item, size=8, color=SOFT_WHITE)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, cost, size=8, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 16 * mm)

# Total
card_flat(c, M + 2 * mm, 228 * mm, PW - 2 * M - 4 * mm, 7 * mm, ACCENT_ORANGE)
txt(c, M + 6 * mm, 230 * mm, "Total COGS", size=10, color=WHITE_HEX, bold=True)
txt(c, M + 6 * mm, 230 * mm, "~$1.14  (50k: ~$0.92)", size=10, color=WHITE_HEX, bold=True,
    align="right", max_w=PW - 2 * M - 16 * mm)

footer(c, 6, TOTAL_PAGES)
c.showPage()

# ── PAGE 7: Hardest Problems ─────────────────────────────────
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "Hardest Problems", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

hard_problems = [
    ("Dynamic seal at the split head joint",
     "The moving bristle half oscillates ~1-2mm at 150+ Hz while "
     "the TPE boot must keep water out of the motor cavity. This "
     "is a fatigue + sealing problem -- the boot must survive "
     "millions of flex cycles in a wet, toothpaste-laden environment "
     "without cracking or delaminating from the PP substrate."),
    ("Insert-molding bristle seal quality at speed",
     "Each brush head has ~35 bristle tufts penetrating the PP "
     "surface. At production speed (seconds per unit), every tuft "
     "base must be fully sealed by the injection process. A single "
     "unsealed tuft is a water ingress path. Narrow process window "
     "for injection pressure, temperature, and bristle positioning."),
    ("BOM discipline with dual-injection process",
     "Overmolding adds ~$0.10-0.20 per unit vs. single-shot molding, "
     "and tooling is 2-3x higher ($15k-25k). With $3-5 retail and "
     "distribution margins to cover, every component must be "
     "ruthlessly cost-optimized. No room for any extras."),
]

for i, (title, desc) in enumerate(hard_problems):
    y_top = 30 * mm + i * 62 * mm
    card(c, M, y_top, PW - 2 * M, 56 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, ACCENT_RED)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=13, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 18 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

card_flat(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "All three problems are manufacturing challenges, not design unknowns. "
         "The Oral-B Pulsar proves the architecture works -- the risk is execution at lower cost.",
         size=11, color=ACCENT_ORANGE, max_w=PW - 2 * M - 8 * mm, line_h=14)

footer(c, 7, TOTAL_PAGES)
c.showPage()

# ── PAGE 8: Gate Result & Next ────────────────────────────────
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "Gate Result & Next", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

# Gate badge
card(c, M, 28 * mm, PW - 2 * M, 22 * mm, CARD_BG)
card(c, M + 4 * mm, 31 * mm, 50 * mm, 16 * mm, ACCENT_GREEN)
txt(c, M + 8 * mm, 35 * mm, "GATE: PASS", size=16, color=WHITE_HEX, bold=True)
txt(c, M + 8 * mm, 43 * mm, "40 pass / 47 N/A / 2 fail", size=9, color=WHITE_HEX)
txt(c, M + 60 * mm, 36 * mm, "System description complete.", size=11, color=WHITE_HEX)
txt(c, M + 60 * mm, 44 * mm, "Ready for mold design.", size=11, color=ACCENT_GREEN, bold=True)

# Power summary
card(c, M, 54 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 58 * mm, "POWER", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 66 * mm, "80 mA active  |  < 1 uA off  |  AAA 1000 mAh",
    size=11, color=WHITE_HEX)
txt(c, M + 4 * mm, 73 * mm, "~750 min runtime  |  ~187 days @ 4 min/day  |  non-rechargeable",
    size=10, color=LIGHT_GRAY)

# Key specs
card(c, M, 80 * mm, PW - 2 * M, 22 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 84 * mm, "KEY SPECS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 92 * mm, "~160mm long  |  ~15mm dia  |  IPX5  |  Single SKU",
    size=10, color=WHITE_HEX)
txt(c, M + 4 * mm, 99 * mm, "No MCU  |  No PCB  |  No app  |  $3-5 retail",
    size=10, color=LIGHT_GRAY)

# Open items
bar(c, M, 108 * mm, PW - 2 * M, 1 * mm, ACCENT_ORANGE)
txt(c, M, 112 * mm, "7 OPEN ITEMS", size=11, color=ACCENT_ORANGE, bold=True)

open_items = [
    ("M1", "Motor performance validation at end-of-life voltage (~1.0V)"),
    ("M2", "Living hinge geometry: thickness, width, PP grade for 4M cycles"),
    ("M2", "TPE head boot fatigue life in wet toothpaste environment"),
    ("M3", "Insert-mold bristle seal QC: tuft pull-force per ISO 20126"),
    ("M3", "Retail channel acceptance test (shelf appeal, blister design)"),
    ("M4", "O-ring compression verification for production QC"),
    ("M4", "Motor current draw screening for dead/high-draw units"),
]

for i, (milestone, desc) in enumerate(open_items):
    y_top = 118 * mm + i * 12 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 10 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3 * mm, milestone, size=9, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 18 * mm, y_top + 3 * mm, desc, size=9, color=SOFT_WHITE)

# CTA
card(c, M, 210 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 214 * mm,
         "Next step: source motor samples and AAA cells, build a "
         "hand-assembled prototype in a 3D-printed shell, and validate "
         "the split-head oscillation and motor stall voltage. Then "
         "commit to dual-injection tooling.",
         size=10, color=WHITE_HEX, max_w=PW - 2 * M - 8 * mm, line_h=13)

footer(c, 8, TOTAL_PAGES)
c.showPage()

c.save()
print(f"Saved {TOTAL_PAGES}-page carousel PDF to {output}")

# ═══════════════════════════════════════════════════════════════
# PPTX GENERATION
# ═══════════════════════════════════════════════════════════════
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(9.375)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_CARD_BG = RGBColor(0x22, 0x22, 0x3A)
C_CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)
C_ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
C_ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
C_ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
C_ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

PM = Inches(0.4)


def pptx_bg(slide, color=C_DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def pptx_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_text(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return txBox


def pptx_add_para(tf, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def pptx_footer(slide, page, total):
    pptx_text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.3),
              Inches(0.6), Inches(0.25), f"{page}/{total}",
              size=8, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ── PPTX PAGE 1: Title ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.5), SLIDE_W - 2 * PM, Inches(0.6),
          "Consumable Electric", size=34, bold=True)
pptx_text(slide, PM, Inches(1.0), SLIDE_W - 2 * PM, Inches(0.6),
          "Toothbrush", size=34, bold=True)

pptx_rect(slide, PM, Inches(1.5), Inches(1.4), Inches(0.04), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(1.7), SLIDE_W - 2 * PM, Inches(0.4),
          "Use it. Toss it. Open a new one.", size=18, color=C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(2.2), SLIDE_W - 2 * PM, Inches(0.6),
          "A battery-powered vibrating toothbrush designed as a true "
          "consumable. AAA alkaline, no MCU, dual-injection sealed. "
          "$3-5 retail vs. Oral-B Pulsar at $8-12.",
          size=11, color=C_LIGHT_GRAY)

img_path = os.path.join(_DIR, "cross_section_illustration_consumable_toothbrush.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, PM, Inches(3.0),
                             width=SLIDE_W - 2 * PM, height=Inches(5.3))

pptx_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_CARD_BG)
pptx_text(slide, PM, SLIDE_H - Inches(0.3), SLIDE_W - 2 * PM, Inches(0.2),
          "Product Overview  |  Concept Stage  |  2026", size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 1, TOTAL_PAGES)

# ── PPTX PAGE 2: The Problem ─────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "The Problem", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_problems = [
    ("Oral-B Pulsar is overpriced for a disposable",
     "At $8-12, the Pulsar pretends to be reusable but has a sealed, "
     "non-replaceable battery. When it dies, you throw it away.",
     C_ACCENT_RED),
    ("Manual brushes don't clean as well",
     "Powered vibration removes more plaque than manual brushing. But "
     "rechargeable brushes cost $30-200. Budget consumers are stuck.",
     C_ACCENT_ORANGE),
    ("No honest sub-$5 powered option exists",
     "Premium rechargeable ($30+) or Oral-B Pulsar ($8-12). Nothing "
     "below $5 where powered brushing becomes truly disposable.",
     C_ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(pptx_problems):
    y = Inches(1.15 + i * 2.15)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
    pptx_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.06), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.3),
              title, size=14, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.65), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
              desc, size=11, color=C_LIGHT_GRAY)

pptx_rounded_rect(slide, PM, Inches(7.7), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
          "TARGET USERS", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.1), SLIDE_W - 2 * PM, Inches(0.25),
          "Budget-conscious consumers  |  Travelers", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Hotels & hospitality  |  Multi-pack buyers", size=12, bold=True)

pptx_footer(slide, 2, TOTAL_PAGES)

# ── PPTX PAGE 3: How It Works ────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "How It Works", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

pptx_steps = [
    ("1. Open", "Tear open blister pack",
     "AAA alkaline battery is pre-installed with years of shelf life. "
     "No charging, no pairing, no setup.",
     C_ACCENT_BLUE),
    ("2. Press", "Push the button once",
     "Latching push-button turns motor on. ERM vibration motor spins "
     "up instantly. Bristles oscillate via split-head mechanism.",
     C_ACCENT_GREEN),
    ("3. Brush", "2 minutes, twice daily",
     "Split bristle head concentrates oscillation at bristle tips. "
     "Rinse under tap after use. IPX5 sealed.",
     C_ACCENT_PURPLE),
    ("4. Toss", "~90 days later, motor stops",
     "Alkaline cell drops below stall voltage. Clear end-of-life: "
     "it works or it doesn't. Discard and open a new one.",
     C_ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(pptx_steps):
    y = Inches(1.1 + i * 1.95)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.35),
              title, size=18, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.45), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              subtitle, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.8), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_footer(slide, 3, TOTAL_PAGES)

# ── PPTX PAGE 4: Architecture ────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Architecture", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_BLUE)

# Power chain
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(1.15), SLIDE_W - 2 * PM, Inches(0.2),
          "Power Path", size=11, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.45), SLIDE_W - 2 * PM, Inches(0.25),
          "AAA Cell --> Latching Switch --> ERM Motor", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.8), SLIDE_W - 2 * PM, Inches(0.2),
          "1.5V alkaline      push on/off       60-100mA, direct drive",
          size=8, color=C_LIGHT_GRAY)

# Mechanical chain
pptx_rounded_rect(slide, PM, Inches(2.35), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Mechanical Path", size=11, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.7), SLIDE_W - 2 * PM, Inches(0.25),
          "Motor --> Eccentric --> Linkage --> Split Head --> Bristles",
          size=10, bold=True)

pptx_rect(slide, PM, Inches(3.4), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_BLUE)
pptx_text(slide, PM, Inches(3.5), SLIDE_W - 2 * PM, Inches(0.2),
          "SUBSYSTEMS", size=11, color=C_ACCENT_BLUE, bold=True)

pptx_subs = [
    ("AAA Alkaline Cell (1.5V)", "1000-1200 mAh, ~750 min runtime, years of shelf life", C_ACCENT_ORANGE),
    ("Latching Push-Button", "Mechanical push-on/push-off, sealed by TPE overmold", C_ACCENT_GREEN),
    ("Cylindrical ERM Motor", "6x12mm, 60-100mA at 1.5V, axial mount in handle neck", C_ACCENT_PURPLE),
    ("Split Bristle Head", "Fixed + moving halves, living hinge pivot, 1-2mm oscillation", C_ACCENT_GREEN),
    ("Rigid PP Body (1st shot)", "Handle, motor pocket, battery tube, insert-molded bristles", C_ACCENT_BLUE),
    ("TPE Overmold (2nd shot)", "Seals all penetrations, grip, button membrane, head boot", C_ACCENT_ORANGE),
    ("Battery Cap + O-ring", "Threaded PP cap, static radial seal, factory-installed", C_ACCENT_PURPLE),
]

for i, (name, desc, color) in enumerate(pptx_subs):
    y = Inches(3.8 + i * 0.58)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.48), bg_c)
    pptx_rect(slide, PM, y, Inches(0.06), Inches(0.48), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), SLIDE_W - 2 * PM, Inches(0.2),
              name, size=10, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.25), SLIDE_W - 2 * PM, Inches(0.2),
              desc, size=8, color=C_LIGHT_GRAY)

# No software callout
pptx_rounded_rect(slide, PM, Inches(7.9), SLIDE_W - 2 * PM, Inches(0.95), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(7.95), SLIDE_W - 2 * PM, Inches(0.25),
          "Zero software. Zero electronics.", size=13, color=C_ACCENT_RED, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.3), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "No MCU, no PCB, no firmware, no app, no cloud. The entire "
          "electrical system is 3 components and 2 wires.",
          size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 4, TOTAL_PAGES)

# ── PPTX PAGE 5: Dual-Injection Sealing ──────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Dual-Injection Sealing", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.05), SLIDE_W - 2 * PM, Inches(0.4),
          "How do you waterproof a $1.50 product for 90 days of wet use?",
          size=13, color=C_LIGHT_GRAY)

# 1st shot
pptx_rounded_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(1.8), C_CARD_BG)
pptx_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_BLUE)
pptx_text(slide, PM + Inches(0.2), Inches(1.75), SLIDE_W - 2 * PM, Inches(0.3),
          "1st Shot: Rigid PP Body", size=16, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(2.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
          "Injection-mold the structural handle with bristles pre-loaded "
          "in the mold cavity. Plastic flows around each bristle base, "
          "creating a watertight seal at every tuft. No secondary sealing.",
          size=11)

# 2nd shot
pptx_rounded_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(2.0), C_CARD_BG)
pptx_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.2), Inches(3.75), SLIDE_W - 2 * PM, Inches(0.3),
          "2nd Shot: TPE Overmold", size=16, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(4.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.2),
          "After motor, wiring, switch, and battery are assembled into "
          "the rigid body, the 2nd injection encapsulates everything. "
          "TPE bonds chemically to PP, sealing all penetrations in one "
          "step: grip texture, button membrane, and dynamic head boot.",
          size=11)

# 4 seal zones
pptx_rounded_rect(slide, PM, Inches(5.9), SLIDE_W - 2 * PM, Inches(2.9), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.2), Inches(6.0), SLIDE_W - 2 * PM, Inches(0.25),
          "4 Seal Zones", size=13, color=C_ACCENT_ORANGE, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(6.4), SLIDE_W - 2 * PM - Inches(0.4), Inches(2.2),
                  "1. Bristle insert-mold: PP flows around nylon tuft bases during 1st injection",
                  size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
pptx_add_para(tf, "2. TPE body overmold: 2nd shot seals handle, switch, wire penetrations",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "3. TPE head boot: Flexible boot around split-head pivot, survives millions of flex cycles",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "4. O-ring battery cap: Threaded PP cap compresses O-ring at handle base",
              size=9, color=C_LIGHT_GRAY, space_before=6)

pptx_footer(slide, 5, TOTAL_PAGES)

# ── PPTX PAGE 6: Constraints & BOM ───────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Constraints & BOM", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_ORANGE)

pptx_cons = [
    ("BOM cost", "< $1.50", "At 10k units; < $1.20 at 50k"),
    ("Motor runtime", ">= 350 min", "~750 min calculated (AAA @ 80mA)"),
    ("Water resistance", "IPX5", "Dual-injection + O-ring, 90 days wet use"),
    ("Handle diameter", "<= 16mm", "AAA cell (10.5mm) + PP wall + TPE"),
    ("Retail price", "$3 - $5", "Undercutting Oral-B Pulsar ($8-12)"),
    ("Shelf life", "> 2 years", "Alkaline chemistry, sealed blister pack"),
]

for i, (name, value, note) in enumerate(pptx_cons):
    y = Inches(1.1 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), Inches(3), Inches(0.2),
              name, size=11, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.3), Inches(5), Inches(0.2),
              note, size=8, color=C_LIGHT_GRAY)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.2),
              value, size=11, bold=True, align=PP_ALIGN.RIGHT)

# BOM
pptx_rect(slide, PM, Inches(5.15), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(5.25), SLIDE_W - 2 * PM, Inches(0.2),
          "BOM ESTIMATE (10k units)", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_bom = [
    ("ERM motor (cylindrical, 6x12mm)", "$0.12"),
    ("AAA alkaline cell", "$0.07"),
    ("Latching push-button switch", "$0.03"),
    ("Wiring (2x, 26AWG, tinned)", "$0.01"),
    ("Spring contact + plate contact", "$0.02"),
    ("PA-612 nylon bristle tufts (x35)", "$0.02"),
    ("O-ring (battery cap seal)", "$0.01"),
    ("PP resin (1st shot, ~8g)", "$0.02"),
    ("TPE resin (2nd shot, ~3g)", "$0.02"),
    ("Battery cap (threaded PP)", "$0.02"),
    ("Blister packaging + card", "$0.08"),
    ("Assembly + test (90-100s/unit)", "$0.72"),
]

for i, (item, cost) in enumerate(pptx_bom):
    y = Inches(5.55 + i * 0.27)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rect(slide, PM + Inches(0.08), y, SLIDE_W - 2 * PM - Inches(0.16), Inches(0.23), bg_c)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), Inches(4), Inches(0.18),
              item, size=8, color=C_SOFT_WHITE)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.18),
              cost, size=8, bold=True, align=PP_ALIGN.RIGHT)

# Total
pptx_rect(slide, PM + Inches(0.08), Inches(8.82), SLIDE_W - 2 * PM - Inches(0.16), Inches(0.3), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), Inches(2), Inches(0.25),
          "Total COGS", size=10, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.25),
          "~$1.14  (50k: ~$0.92)", size=10, bold=True, align=PP_ALIGN.RIGHT)

pptx_footer(slide, 6, TOTAL_PAGES)

# ── PPTX PAGE 7: Hardest Problems ────────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Hardest Problems", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_hard = [
    ("Dynamic seal at the split head joint",
     "The moving bristle half oscillates ~1-2mm at 150+ Hz while "
     "the TPE boot must keep water out of the motor cavity. Fatigue + "
     "sealing problem -- the boot must survive millions of flex cycles "
     "in a wet, toothpaste-laden environment."),
    ("Insert-molding bristle seal quality at speed",
     "~35 bristle tufts penetrate the PP surface. At production speed, "
     "every tuft base must be sealed by injection. A single unsealed "
     "tuft is a water ingress path. Narrow process window."),
    ("BOM discipline with dual-injection process",
     "Overmolding adds ~$0.10-0.20/unit vs. single-shot, tooling is "
     "2-3x higher ($15k-25k). With $3-5 retail and distribution "
     "margins, every component must be ruthlessly cost-optimized."),
]

for i, (title, desc) in enumerate(pptx_hard):
    y = Inches(1.15 + i * 2.4)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.15), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              f"{i+1}. {title}", size=13, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.55), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.4),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_rounded_rect(slide, PM, Inches(8.4), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "All three are manufacturing challenges, not design unknowns. "
          "Oral-B Pulsar proves the architecture -- the risk is execution at lower cost.",
          size=11, color=C_ACCENT_ORANGE)

pptx_footer(slide, 7, TOTAL_PAGES)

# ── PPTX PAGE 8: Gate Result & Next ──────────────────────────
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Gate Result & Next", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

# Gate badge
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_rect(slide, PM + Inches(0.15), Inches(1.2), Inches(2.0), Inches(0.6), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.3), Inches(1.22), Inches(1.8), Inches(0.3),
          "GATE: PASS", size=16, bold=True)
pptx_text(slide, PM + Inches(0.3), Inches(1.52), Inches(1.8), Inches(0.2),
          "40 pass / 47 N/A / 2 fail", size=9)
pptx_text(slide, PM + Inches(2.4), Inches(1.28), Inches(4), Inches(0.2),
          "System description complete.", size=11)
pptx_text(slide, PM + Inches(2.4), Inches(1.55), Inches(4), Inches(0.2),
          "Ready for mold design.", size=11, color=C_ACCENT_GREEN, bold=True)

# Power summary
pptx_rounded_rect(slide, PM, Inches(2.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.15), SLIDE_W - 2 * PM, Inches(0.2),
          "POWER", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "80 mA active  |  < 1 uA off  |  AAA 1000 mAh", size=11)
pptx_text(slide, PM + Inches(0.15), Inches(2.65), SLIDE_W - 2 * PM, Inches(0.2),
          "~750 min runtime  |  ~187 days @ 4 min/day  |  non-rechargeable",
          size=10, color=C_LIGHT_GRAY)

# Key specs
pptx_rounded_rect(slide, PM, Inches(3.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(3.15), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY SPECS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(3.4), SLIDE_W - 2 * PM, Inches(0.2),
          "~160mm long  |  ~15mm dia  |  IPX5  |  Single SKU", size=10)
pptx_text(slide, PM + Inches(0.15), Inches(3.65), SLIDE_W - 2 * PM, Inches(0.2),
          "No MCU  |  No PCB  |  No app  |  $3-5 retail",
          size=10, color=C_LIGHT_GRAY)

# Open items
pptx_rect(slide, PM, Inches(4.2), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_ORANGE)
pptx_text(slide, PM, Inches(4.3), SLIDE_W - 2 * PM, Inches(0.2),
          "7 OPEN ITEMS", size=11, color=C_ACCENT_ORANGE, bold=True)

pptx_open = [
    ("M1", "Motor performance validation at end-of-life voltage (~1.0V)"),
    ("M2", "Living hinge geometry: thickness, width, PP grade for 4M cycles"),
    ("M2", "TPE head boot fatigue life in wet toothpaste environment"),
    ("M3", "Insert-mold bristle seal QC: tuft pull-force per ISO 20126"),
    ("M3", "Retail channel acceptance test (shelf appeal, blister design)"),
    ("M4", "O-ring compression verification for production QC"),
    ("M4", "Motor current draw screening for dead/high-draw units"),
]

for i, (milestone, desc) in enumerate(pptx_open):
    y = Inches(4.6 + i * 0.47)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.38), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.07), Inches(0.5), Inches(0.2),
              milestone, size=9, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.7), y + Inches(0.07), SLIDE_W - 2 * PM - Inches(1.0), Inches(0.2),
              desc, size=9, color=C_SOFT_WHITE)

# CTA
pptx_rounded_rect(slide, PM, Inches(8.2), SLIDE_W - 2 * PM, Inches(0.8), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.25), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.7),
          "Next step: source motor samples and AAA cells, build a "
          "hand-assembled prototype in a 3D-printed shell, and validate "
          "split-head oscillation and motor stall voltage. Then commit "
          "to dual-injection tooling.",
          size=10)

pptx_footer(slide, 8, TOTAL_PAGES)

# Save PPTX
pptx_output = os.path.join(_DIR, "Consumable_Toothbrush_Carousel.pptx")
prs.save(pptx_output)
print(f"Saved {TOTAL_PAGES}-slide carousel PPTX to {pptx_output}")
