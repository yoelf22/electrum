#!/usr/bin/env python3
"""Build a LinkedIn carousel PDF and PPTX for Shusher.

Format: 1080x1350 px (4:5 portrait) — optimized for mobile feed.
Uses reportlab for PDF and python-pptx for PPTX generation.
"""

import os
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor, white
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Page size: 4:5 ratio --
PW = 190 * mm
PH = 237.5 * mm

# -- Colors --
DARK_BG = "#1A1A2E"
CARD_BG = "#22223A"
CARD_BG_ALT = "#1E1E34"
ACCENT_ORANGE = "#FF8C00"
ACCENT_GREEN = "#4EC978"
ACCENT_RED = "#FF4545"
ACCENT_BLUE = "#009BF5"
ACCENT_PURPLE = "#A855F7"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY = "#BBBBCC"
SOFT_WHITE = "#F0F0F5"

M = 10 * mm  # standard margin


def bg(c, color=DARK_BG):
    c.setFillColor(HexColor(color))
    c.rect(0, 0, PW, PH, fill=1, stroke=0)


def accent_strip(c, color=ACCENT_ORANGE):
    c.setFillColor(HexColor(color))
    c.rect(0, PH - 2 * mm, PW, 2 * mm, fill=1, stroke=0)


def card(c, x, y, w, h, color=CARD_BG):
    """Draw a card. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    c.roundRect(x, PH - y - h, w, h, 3 * mm, fill=1, stroke=0)


def card_flat(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def txt(c, x, y, text, size=14, color=WHITE_HEX, bold=False, align="left", max_w=None):
    """Draw text. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    ry = PH - y - size * 0.35
    if align == "center" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + (max_w - tw) / 2, ry, text)
    elif align == "right" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + max_w - tw, ry, text)
    else:
        c.drawString(x, ry, text)


def txt_wrap(c, x, y, text, size=11, color=WHITE_HEX, bold=False, line_h=None, max_w=None):
    """Draw wrapped text. Returns y after last line."""
    if line_h is None:
        line_h = size * 1.4
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFillColor(HexColor(color))
    c.setFont(font, size)
    if max_w is None:
        max_w = PW - 2 * M
    words = text.split()
    lines = []
    current = ""
    for w in words:
        test = current + (" " if current else "") + w
        if c.stringWidth(test, font, size) > max_w:
            if current:
                lines.append(current)
            current = w
        else:
            current = test
    if current:
        lines.append(current)
    for line in lines:
        c.drawString(x, PH - y - size * 0.35, line)
        y += line_h
    return y


def bar(c, x, y, w, h, color):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def circle_num(c, x, y, num, color):
    r = 4 * mm
    c.setFillColor(HexColor(color))
    c.circle(x + r, PH - y - r, r, fill=1, stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    tw = c.stringWidth(str(num), "Helvetica-Bold", 12)
    c.drawString(x + r - tw / 2, PH - y - r - 4, str(num))


def footer(c, page, total):
    c.setFillColor(HexColor(LIGHT_GRAY))
    c.setFont("Helvetica", 8)
    label = f"{page}/{total}"
    tw = c.stringWidth(label, "Helvetica", 8)
    c.drawString(PW - M - tw, 4 * mm, label)


TOTAL_PAGES = 8
output = os.path.join(_DIR, "Shusher_Carousel.pdf")
c = canvas.Canvas(output, pagesize=(PW, PH))

# ================================================================
# PAGE 1: Title
# ================================================================
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 14 * mm, "Shusher", size=42, color=WHITE_HEX, bold=True)

bar(c, M, 28 * mm, 35 * mm, 1 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 34 * mm,
         "The polite nudge you never had to give.",
         size=18, color=ACCENT_GREEN, max_w=PW - 2 * M)

txt_wrap(c, M, 48 * mm,
         "A palm-sized puck that listens for sustained loud speech, "
         "then fires a discreet, human-sounding \"shhh\" toward the "
         "offender. You aim it. It does the rest.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Image
img_path = os.path.join(_DIR, "cross_section_illustration_shusher.png")
if os.path.exists(img_path):
    img_top_y = 70 * mm
    img_max_w = PW - 2 * M
    img_max_h = 140 * mm
    c.drawImage(ImageReader(img_path),
                M, PH - img_top_y - img_max_h,
                width=img_max_w, height=img_max_h,
                preserveAspectRatio=True, anchor='sw', mask='auto')

# Bottom bar
card_flat(c, 0, PH - 8 * mm, PW, 8 * mm, CARD_BG)
txt(c, M, PH - 5 * mm, "Product Overview  |  Concept Stage  |  2026",
    size=9, color=LIGHT_GRAY)

footer(c, 1, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 2: The Problem
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "The Problem", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Loud talkers and speakerphone abusers",
     "Cafes, coworking spaces, trains -- someone always decides "
     "the whole room needs to hear their call. You tolerate it, "
     "or you confront a stranger. Neither feels good.",
     ACCENT_RED),
    ("Confrontation is awkward",
     "Asking someone to lower their voice means making yourself "
     "the center of attention. Most people avoid it. The loud "
     "talker wins by default.",
     ACCENT_ORANGE),
    ("No passive, low-effort solution exists",
     "Noise-cancelling headphones help you, but they don't signal "
     "anything to the offender. There's no device that delivers a "
     "social nudge without making you the messenger.",
     ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    y_top = 30 * mm + i * 55 * mm
    card(c, M, y_top, PW - 2 * M, 49 * mm, CARD_BG)
    bar(c, M, y_top, PW - 2 * M, 1.5 * mm, color)
    txt(c, M + 5 * mm, y_top + 8 * mm, title, size=16, color=color, bold=True)
    txt_wrap(c, M + 5 * mm, y_top + 20 * mm, desc,
             size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Target users
card_flat(c, M, 200 * mm, PW - 2 * M, 30 * mm, CARD_BG)
txt(c, M + 4 * mm, 204 * mm, "TARGET USERS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 213 * mm, "Remote workers  |  Students  |  Commuters",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 222 * mm, "Library goers  |  Introverts  |  Anyone who values quiet",
    size=12, color=WHITE_HEX, bold=True)

footer(c, 2, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 3: How It Works
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "How It Works", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

steps = [
    ("Aim", "Point the puck at the noise source",
     "Place Shusher on your table and aim the front face "
     "toward the loud talker. Press the side button once. "
     "A brief vibration confirms it's active.",
     ACCENT_BLUE),
    ("Detect", "Device listens and compares",
     "Two mics -- front (cardioid) and rear (omni) -- "
     "measure audio levels continuously. When the front mic "
     "detects sustained loud speech 10+ dB above the ambient "
     "baseline for 6+ seconds, the trigger fires.",
     ACCENT_GREEN),
    ("Shush", "A natural \"shhh\" fires forward",
     "A pre-recorded human shush plays through a front-facing "
     "speaker with waveguide. Audible at the offender's table "
     "(2-5m), but blends with cafe noise for everyone else. "
     "Sounds human. Plausibly deniable.",
     ACCENT_PURPLE),
    ("Cool", "Wait, then repeat if needed",
     "After a shush, a 30-60 second cooldown prevents "
     "re-triggering. One polite nudge per incident. "
     "If they get loud again, another shush fires.",
     ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(steps):
    y_top = 28 * mm + i * 50 * mm
    card(c, M, y_top, PW - 2 * M, 44 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, color)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=18, color=color, bold=True)
    txt(c, M + 16 * mm, y_top + 16 * mm, subtitle, size=11, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 26 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom note
card_flat(c, M, 232 * mm, PW - 2 * M, 4 * mm, CARD_BG)

footer(c, 3, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 4: Architecture
# ================================================================
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 12 * mm, "Architecture", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_BLUE)

# Detection chain
card(c, M, 28 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 32 * mm, "Detection Chain", size=11, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 41 * mm, "Front mic  -->  Bandpass  -->  RMS  -->  Differential  -->  Trigger",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 49 * mm, "cardioid PDM    300-4kHz IIR    50ms window    front vs rear    sustained >6s",
    size=8, color=LIGHT_GRAY)

# Shush chain
card(c, M, 60 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 64 * mm, "Shush Chain", size=11, color=ACCENT_GREEN, bold=True)
txt(c, M + 4 * mm, 72 * mm, "Flash sample  -->  I2S  -->  MAX98357A  -->  Speaker + Waveguide",
    size=10, color=WHITE_HEX, bold=True)

# Subsystems
bar(c, M, 88 * mm, PW - 2 * M, 1 * mm, ACCENT_BLUE)
txt(c, M, 92 * mm, "SUBSYSTEMS", size=11, color=ACCENT_BLUE, bold=True)

subsystems = [
    ("MCU (ESP32-S3)", "Dual-core 240 MHz, audio DSP, BLE 5.0, 8MB flash", ACCENT_GREEN),
    ("Front mic (ICS-43434)", "Cardioid via acoustic porting, PDM 16 kHz, 65 dB SNR", ACCENT_PURPLE),
    ("Rear mic (ICS-43434)", "Omnidirectional, ambient baseline reference", ACCENT_GREEN),
    ("Speaker (28mm + waveguide)", "Forward-biased shush output, 80-85 dB @ 10cm", ACCENT_BLUE),
    ("Amp (MAX98357A)", "I2S class-D, auto-shutdown when idle", ACCENT_ORANGE),
    ("Power (1200 mAh LiPo)", "USB-C charging, ~3 days per charge", ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(subsystems):
    y_top = 98 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 4 * mm, name, size=11, color=WHITE_HEX, bold=True)
    txt(c, M + 4 * mm, y_top + 11 * mm, desc, size=8, color=LIGHT_GRAY)
    bar(c, M, y_top, 1.5 * mm, 14 * mm, color)

# Discretion model
card(c, M, 202 * mm, PW - 2 * M, 28 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 206 * mm, "Designed to disappear", size=13, color=ACCENT_PURPLE, bold=True)
txt_wrap(c, M + 4 * mm, 216 * mm,
         "No visible LEDs during operation. No electronic sound artifacts. "
         "The shush sounds human -- breathy, variable, naturally timed. "
         "The offender wonders if someone nearby shushed them. Not a gadget.",
         size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 8 * mm, line_h=12)

footer(c, 4, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 5: Plausible Deniability (Key Innovation)
# ================================================================
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 12 * mm, "Plausible Deniability", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 28 * mm,
         "The product works only if nobody knows it's there.",
         size=13, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Looks like a coaster
card(c, M, 42 * mm, PW - 2 * M, 45 * mm, CARD_BG)
bar(c, M, 42 * mm, PW - 2 * M, 1.5 * mm, ACCENT_GREEN)
txt(c, M + 5 * mm, 48 * mm, "Looks like a coaster", size=16, color=ACCENT_GREEN, bold=True)
txt_wrap(c, M + 5 * mm, 60 * mm,
         "70mm diameter, 15mm tall, matte dark puck. No branding on top. "
         "Speaker slot disguised as a seam line. Mic port looks like a "
         "mold mark. Sits naturally on any cafe table.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Sounds human
card(c, M, 94 * mm, PW - 2 * M, 55 * mm, CARD_BG)
bar(c, M, 94 * mm, PW - 2 * M, 1.5 * mm, ACCENT_PURPLE)
txt(c, M + 5 * mm, 100 * mm, "Sounds human", size=16, color=ACCENT_PURPLE, bold=True)
txt_wrap(c, M + 5 * mm, 112 * mm,
         "20+ pre-recorded shush variants across 5 styles: classic \"shhh,\" "
         "throat clear, gentle \"ahem,\" whispered \"excuse me,\" and passive-aggressive "
         "sigh. Each firing picks a different sample. Breathy, slightly variable, "
         "naturally timed. No robotic quality, no beeps, no electronic artifacts.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Disappears after firing
card(c, M, 158 * mm, PW - 2 * M, 70 * mm, CARD_BG_ALT)
txt(c, M + 5 * mm, 163 * mm, "The discretion stack", size=13, color=ACCENT_ORANGE, bold=True)

discretion = [
    ("No LEDs during operation",
     "A single LED under the device confirms power-on, "
     "then goes dark. Invisible on a table."),
    ("Forward-biased audio",
     "Waveguide provides ~4-6 dB front-to-side ratio. "
     "The target hears it; side-facing listeners don't."),
    ("Cooldown prevents harassment",
     "Max one shush per 30-60 seconds. A polite nudge, "
     "not a weapon. The offender self-corrects or doesn't."),
]

ry = 175 * mm
for title, desc in discretion:
    txt(c, M + 5 * mm, ry, f"- {title}:", size=10, color=WHITE_HEX, bold=True)
    ry = txt_wrap(c, M + 8 * mm, ry + 12, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 16 * mm, line_h=12)
    ry += 3 * mm

footer(c, 5, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 6: Constraints & BOM
# ================================================================
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 12 * mm, "Constraints & BOM", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_ORANGE)

constraints = [
    ("Battery life", "~3 days", "1200 mAh LiPo, 4-hr daily active sessions"),
    ("Size", "70mm x 15mm", "Coaster-sized puck, fits naturally on a table"),
    ("Shush range", "2-5 meters", "Audible to the target, blends beyond 5m"),
    ("Trigger latency", "~6 seconds", "Sustained loud speech, not brief bursts"),
    ("False positive rate", "< 1 / hour", "Shushing the barista is a product killer"),
    ("Weight", "~65g", "PCB 20g + battery 25g + enclosure 15g + speaker 5g"),
]

for i, (name, value, note) in enumerate(constraints):
    y_top = 28 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=11, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, note, size=8, color=LIGHT_GRAY)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, value, size=11, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 8 * mm)

# BOM
bar(c, M, 133 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 137 * mm, "BOM ESTIMATE (1k units)", size=11, color=ACCENT_GREEN, bold=True)

bom = [
    ("ESP32-S3-WROOM-1 (N8)", "$2.80"),
    ("Front MEMS mic (ICS-43434)", "$0.80"),
    ("Rear MEMS mic (ICS-43434)", "$0.80"),
    ("MAX98357A class-D amp", "$1.20"),
    ("28mm dynamic speaker", "$0.60"),
    ("LRA haptic motor (8mm)", "$0.50"),
    ("LiPo battery (1200 mAh)", "$2.50"),
    ("MCP73831 charge IC + USB-C", "$0.80"),
    ("LDO + passives + RGB LED", "$0.50"),
    ("PCB (55mm round, 4-layer)", "$1.20"),
    ("Enclosure (2-piece + waveguide)", "$2.50"),
    ("Packaging + assembly + test", "$2.80"),
]

for i, (item, cost) in enumerate(bom):
    y_top = 144 * mm + i * 7 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card_flat(c, M + 2 * mm, y_top, PW - 2 * M - 4 * mm, 6 * mm, bg_c)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, item, size=8, color=SOFT_WHITE)
    txt(c, M + 6 * mm, y_top + 1.5 * mm, cost, size=8, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 16 * mm)

# Total
card_flat(c, M + 2 * mm, 228 * mm, PW - 2 * M - 4 * mm, 7 * mm, ACCENT_ORANGE)
txt(c, M + 6 * mm, 230 * mm, "Total COGS", size=10, color=WHITE_HEX, bold=True)
txt(c, M + 6 * mm, 230 * mm, "~$17.25  (5k: ~$13.50)", size=10, color=WHITE_HEX, bold=True,
    align="right", max_w=PW - 2 * M - 16 * mm)

footer(c, 6, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 7: Hardest Problems
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "Hardest Problems", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

hard_problems = [
    ("Distinguishing shush-worthy noise from cafe ambiance",
     "The front mic picks up everything in the aimed direction -- "
     "the loud talker, but also music, dish clatter, espresso machines. "
     "The classifier must isolate sustained loud speech above the "
     "ambient baseline and hold for seconds before triggering. "
     "Too sensitive: shushes the barista. Too conservative: never fires."),
    ("Making the shush sound convincingly human",
     "If the shush sounds electronic or robotic, the product fails "
     "socially. It needs to sound like a real person -- breathy, slightly "
     "variable, naturally timed. Speaker coloration, waveguide resonances, "
     "and enclosure vibration could add artifacts. "
     "Multiple recordings with randomized selection on each firing."),
    ("Achieving useful directivity from a small speaker",
     "A 70mm puck with a 28mm speaker and short waveguide won't "
     "produce a tight beam. Target: ~4-6 dB front-to-side ratio at "
     "2-4 kHz (the \"shhh\" frequency range). Enough to make the shush "
     "louder for the target than for people at 90 degrees, but not a "
     "laser. Waveguide geometry and speaker placement are critical."),
]

for i, (title, desc) in enumerate(hard_problems):
    y_top = 30 * mm + i * 62 * mm
    card(c, M, y_top, PW - 2 * M, 56 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, ACCENT_RED)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=13, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 18 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom
card_flat(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "All three require physical prototyping. Acoustic porting, "
         "waveguide geometry, and sound recording quality can only be "
         "validated by building and testing in real environments.",
         size=11, color=ACCENT_ORANGE, max_w=PW - 2 * M - 8 * mm, line_h=14)

footer(c, 7, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 8: Gate Result & Open Items
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "Gate Result & Next", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

# Gate badge
card(c, M, 28 * mm, PW - 2 * M, 22 * mm, CARD_BG)
card(c, M + 4 * mm, 31 * mm, 50 * mm, 16 * mm, ACCENT_ORANGE)
txt(c, M + 8 * mm, 35 * mm, "GATE: 74 PASS", size=16, color=WHITE_HEX, bold=True)
txt(c, M + 8 * mm, 43 * mm, "74 pass / 6 N/A / 8 fail", size=9, color=WHITE_HEX)
txt(c, M + 60 * mm, 36 * mm, "System description complete.", size=11, color=WHITE_HEX)
txt(c, M + 60 * mm, 44 * mm, "8 items need prototype validation.", size=11, color=ACCENT_ORANGE, bold=True)

# Power summary
card(c, M, 54 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 58 * mm, "POWER", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 66 * mm, "84 mA active  |  2.3 mA standby  |  1200 mAh battery",
    size=11, color=WHITE_HEX)
txt(c, M + 4 * mm, 73 * mm, "~3 days typical  |  USB-C charge  |  ~2.5 hr 0-100%",
    size=10, color=LIGHT_GRAY)

# Key specs
card(c, M, 80 * mm, PW - 2 * M, 22 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 84 * mm, "KEY SPECS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 92 * mm, "70mm x 15mm puck  |  ~65g  |  ESP32-S3  |  BLE 5.0",
    size=10, color=WHITE_HEX)
txt(c, M + 4 * mm, 99 * mm, "2 MEMS mics  |  Waveguide speaker  |  $50-70 retail",
    size=10, color=LIGHT_GRAY)

# Open items
bar(c, M, 108 * mm, PW - 2 * M, 1 * mm, ACCENT_ORANGE)
txt(c, M, 112 * mm, "KEY OPEN ITEMS", size=11, color=ACCENT_ORANGE, bold=True)

open_items = [
    ("M2", "Cardioid mic porting: 10-15 dB front-to-rear rejection?"),
    ("M2", "Waveguide: >4 dB front-to-side ratio at 2-4 kHz?"),
    ("M2", "Shush naturalness through 28mm speaker in plastic enclosure"),
    ("M3", "False positive rate in 5+ real cafe environments"),
    ("M2", "ESP32-S3 BLE + dual-mic audio coexistence under load"),
    ("M2", "Sound recording session: 20+ shush variants, 5 styles"),
    ("M1", "Market positioning: $50 tool vs. $30 fun gadget"),
    ("M3", "DSP duty-cycle optimization for extended battery life"),
]

for i, (milestone, desc) in enumerate(open_items):
    y_top = 118 * mm + i * 12 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 10 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3 * mm, milestone, size=9, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 18 * mm, y_top + 3 * mm, desc, size=9, color=SOFT_WHITE)

# CTA
card(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "Next step: breadboard prototype on an ESP32-S3 devkit with "
         "two MEMS mics and a small speaker. Test the front/rear "
         "differential in a real cafe. Record the shush samples. "
         "Validate before committing to PCB and enclosure tooling.",
         size=10, color=WHITE_HEX, max_w=PW - 2 * M - 8 * mm, line_h=13)

footer(c, 8, TOTAL_PAGES)
c.showPage()

# ================================================================
c.save()
print(f"Saved {TOTAL_PAGES}-page carousel PDF to {output}")
print(f"Page size: {PW/mm:.0f} x {PH/mm:.0f} mm (4:5 ratio)")

# ================================================================
# PPTX GENERATION
# ================================================================
# 4:5 portrait slide: 7.5 x 9.375 inches (same ratio as PDF)
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(9.375)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# -- PPTX helper colors --
C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_CARD_BG = RGBColor(0x22, 0x22, 0x3A)
C_CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)
C_ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
C_ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
C_ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
C_ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

PM = Inches(0.4)  # margin


def pptx_bg(slide, color=C_DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def pptx_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_text(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return txBox


def pptx_add_para(tf, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def pptx_footer(slide, page, total):
    pptx_text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.3),
              Inches(0.6), Inches(0.25), f"{page}/{total}",
              size=8, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ================================================================
# PPTX PAGE 1: Title
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.5), SLIDE_W - 2 * PM, Inches(0.7),
          "Shusher", size=42, bold=True)

pptx_rect(slide, PM, Inches(1.1), Inches(1.4), Inches(0.04), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.3), SLIDE_W - 2 * PM, Inches(0.4),
          "The polite nudge you never had to give.", size=18, color=C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(1.8), SLIDE_W - 2 * PM, Inches(0.7),
          "A palm-sized puck that listens for sustained loud speech, "
          "then fires a discreet, human-sounding \"shhh\" toward the "
          "offender. You aim it. It does the rest.",
          size=11, color=C_LIGHT_GRAY)

# Image
img_path = os.path.join(_DIR, "cross_section_illustration_shusher.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, PM, Inches(2.8),
                             width=SLIDE_W - 2 * PM, height=Inches(5.5))

pptx_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_CARD_BG)
pptx_text(slide, PM, SLIDE_H - Inches(0.3), SLIDE_W - 2 * PM, Inches(0.2),
          "Product Overview  |  Concept Stage  |  2026", size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 1, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 2: The Problem
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "The Problem", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_problems = [
    ("Loud talkers and speakerphone abusers",
     "Cafes, coworking spaces, trains -- someone always decides "
     "the whole room needs to hear their call. You tolerate it, "
     "or you confront a stranger. Neither feels good.",
     C_ACCENT_RED),
    ("Confrontation is awkward",
     "Asking someone to lower their voice means making yourself "
     "the center of attention. Most people avoid it. The loud "
     "talker wins by default.",
     C_ACCENT_ORANGE),
    ("No passive, low-effort solution exists",
     "Noise-cancelling headphones help you, but they don't signal "
     "anything to the offender. There's no device that delivers a "
     "social nudge without making you the messenger.",
     C_ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(pptx_problems):
    y = Inches(1.15 + i * 2.15)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
    pptx_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.06), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.3),
              title, size=16, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.65), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
              desc, size=11, color=C_LIGHT_GRAY)

# Target users
pptx_rounded_rect(slide, PM, Inches(7.7), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
          "TARGET USERS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.1), SLIDE_W - 2 * PM, Inches(0.25),
          "Remote workers  |  Students  |  Commuters", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Library goers  |  Introverts  |  Anyone who values quiet", size=12, bold=True)

pptx_footer(slide, 2, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 3: How It Works
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "How It Works", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

pptx_steps = [
    ("1. Aim", "Point the puck at the noise source",
     "Place Shusher on your table and aim the front face "
     "toward the loud talker. Press the side button once. "
     "A brief vibration confirms it's active.",
     C_ACCENT_BLUE),
    ("2. Detect", "Device listens and compares",
     "Two mics -- front (cardioid) and rear (omni) -- "
     "measure audio levels continuously. When the front mic "
     "detects sustained loud speech 10+ dB above ambient "
     "for 6+ seconds, the trigger fires.",
     C_ACCENT_GREEN),
    ("3. Shush", "A natural \"shhh\" fires forward",
     "A pre-recorded human shush plays through a front-facing "
     "speaker with waveguide. Audible at 2-5m, blends with "
     "cafe noise for everyone else. Plausibly deniable.",
     C_ACCENT_PURPLE),
    ("4. Cool", "Wait, then repeat if needed",
     "After a shush, a 30-60 second cooldown prevents "
     "re-triggering. One polite nudge per incident. "
     "If they get loud again, another shush fires.",
     C_ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(pptx_steps):
    y = Inches(1.1 + i * 1.95)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.35),
              title, size=18, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.45), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              subtitle, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.8), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_footer(slide, 3, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 4: Architecture
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Architecture", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_BLUE)

# Detection chain
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(1.15), SLIDE_W - 2 * PM, Inches(0.2),
          "Detection Chain", size=11, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Front mic  -->  Bandpass  -->  RMS  -->  Differential  -->  Trigger",
          size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.8), SLIDE_W - 2 * PM, Inches(0.2),
          "cardioid PDM    300-4kHz IIR    50ms window    front vs rear    sustained >6s",
          size=8, color=C_LIGHT_GRAY)

# Shush chain
pptx_rounded_rect(slide, PM, Inches(2.35), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Shush Chain", size=11, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.7), SLIDE_W - 2 * PM, Inches(0.25),
          "Flash sample  -->  I2S  -->  MAX98357A  -->  Speaker + Waveguide",
          size=10, bold=True)

pptx_rect(slide, PM, Inches(3.4), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_BLUE)
pptx_text(slide, PM, Inches(3.5), SLIDE_W - 2 * PM, Inches(0.2),
          "SUBSYSTEMS", size=11, color=C_ACCENT_BLUE, bold=True)

pptx_subsystems = [
    ("MCU (ESP32-S3)", "Dual-core 240 MHz, audio DSP, BLE 5.0, 8MB flash", C_ACCENT_GREEN),
    ("Front mic (ICS-43434)", "Cardioid via acoustic porting, PDM 16 kHz, 65 dB SNR", C_ACCENT_PURPLE),
    ("Rear mic (ICS-43434)", "Omnidirectional, ambient baseline reference", C_ACCENT_GREEN),
    ("Speaker (28mm + waveguide)", "Forward-biased shush output, 80-85 dB @ 10cm", C_ACCENT_BLUE),
    ("Amp (MAX98357A)", "I2S class-D, auto-shutdown when idle", C_ACCENT_ORANGE),
    ("Power (1200 mAh LiPo)", "USB-C charging, ~3 days per charge", C_ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(pptx_subsystems):
    y = Inches(3.8 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_rect(slide, PM, y, Inches(0.06), Inches(0.55), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), SLIDE_W - 2 * PM, Inches(0.2),
              name, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.28), SLIDE_W - 2 * PM, Inches(0.2),
              desc, size=8, color=C_LIGHT_GRAY)

# Discretion model
pptx_rounded_rect(slide, PM, Inches(7.75), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.25),
          "Designed to disappear", size=13, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.15), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.6),
          "No visible LEDs during operation. No electronic sound artifacts. "
          "The shush sounds human -- breathy, variable, naturally timed. "
          "The offender wonders if someone nearby shushed them. Not a gadget.",
          size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 4, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 5: Plausible Deniability
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Plausible Deniability", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.05), SLIDE_W - 2 * PM, Inches(0.4),
          "The product works only if nobody knows it's there.",
          size=13, color=C_LIGHT_GRAY)

# Looks like a coaster
pptx_rounded_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(1.8), C_CARD_BG)
pptx_rect(slide, PM, Inches(1.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.2), Inches(1.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Looks like a coaster", size=16, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(2.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
          "70mm diameter, 15mm tall, matte dark puck. No branding on top. "
          "Speaker slot disguised as a seam line. Mic port looks like a "
          "mold mark. Sits naturally on any cafe table.",
          size=11, color=C_LIGHT_GRAY)

# Sounds human
pptx_rounded_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
pptx_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_PURPLE)
pptx_text(slide, PM + Inches(0.2), Inches(3.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Sounds human", size=16, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(4.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.2),
          "20+ pre-recorded shush variants across 5 styles: classic \"shhh,\" "
          "throat clear, gentle \"ahem,\" whispered \"excuse me,\" and passive-aggressive "
          "sigh. Each firing picks a different sample. Breathy, slightly variable, "
          "naturally timed. No robotic quality, no beeps, no electronic artifacts.",
          size=11, color=C_LIGHT_GRAY)

# Discretion stack
pptx_rounded_rect(slide, PM, Inches(6.1), SLIDE_W - 2 * PM, Inches(2.75), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.2), Inches(6.2), SLIDE_W - 2 * PM, Inches(0.25),
          "The discretion stack", size=13, color=C_ACCENT_ORANGE, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(6.6), SLIDE_W - 2 * PM - Inches(0.4), Inches(2.0),
                  "No LEDs during operation: A single LED under the device confirms "
                  "power-on, then goes dark. Invisible on a table.", size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
pptx_add_para(tf, "Forward-biased audio: Waveguide provides ~4-6 dB front-to-side ratio. "
              "The target hears it; side-facing listeners don't.",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Cooldown prevents harassment: Max one shush per 30-60 seconds. "
              "A polite nudge, not a weapon.",
              size=9, color=C_LIGHT_GRAY, space_before=6)

pptx_footer(slide, 5, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 6: Constraints & BOM
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Constraints & BOM", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_ORANGE)

pptx_constraints = [
    ("Battery life", "~3 days", "1200 mAh LiPo, 4-hr daily active sessions"),
    ("Size", "70mm x 15mm", "Coaster-sized puck, fits naturally on a table"),
    ("Shush range", "2-5 meters", "Audible to the target, blends beyond 5m"),
    ("Trigger latency", "~6 seconds", "Sustained loud speech, not brief bursts"),
    ("False positive rate", "< 1 / hour", "Shushing the barista is a product killer"),
    ("Weight", "~65g", "PCB 20g + battery 25g + enclosure 15g + speaker 5g"),
]

for i, (name, value, note) in enumerate(pptx_constraints):
    y = Inches(1.1 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), Inches(3), Inches(0.2),
              name, size=11, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.3), Inches(5), Inches(0.2),
              note, size=8, color=C_LIGHT_GRAY)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.2),
              value, size=11, bold=True, align=PP_ALIGN.RIGHT)

# BOM
pptx_rect(slide, PM, Inches(5.15), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(5.25), SLIDE_W - 2 * PM, Inches(0.2),
          "BOM ESTIMATE (1k units)", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_bom = [
    ("ESP32-S3-WROOM-1 (N8)", "$2.80"),
    ("Front MEMS mic (ICS-43434)", "$0.80"),
    ("Rear MEMS mic (ICS-43434)", "$0.80"),
    ("MAX98357A class-D amp", "$1.20"),
    ("28mm dynamic speaker", "$0.60"),
    ("LRA haptic motor (8mm)", "$0.50"),
    ("LiPo battery (1200 mAh)", "$2.50"),
    ("MCP73831 charge IC + USB-C", "$0.80"),
    ("LDO + passives + RGB LED", "$0.50"),
    ("PCB (55mm round, 4-layer)", "$1.20"),
    ("Enclosure (2-piece + waveguide)", "$2.50"),
    ("Packaging + assembly + test", "$2.80"),
]

for i, (item, cost) in enumerate(pptx_bom):
    y = Inches(5.55 + i * 0.27)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rect(slide, PM + Inches(0.08), y, SLIDE_W - 2 * PM - Inches(0.16), Inches(0.23), bg_c)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), Inches(4), Inches(0.18),
              item, size=8, color=C_SOFT_WHITE)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.02), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.18),
              cost, size=8, bold=True, align=PP_ALIGN.RIGHT)

# Total
pptx_rect(slide, PM + Inches(0.08), Inches(8.82), SLIDE_W - 2 * PM - Inches(0.16), Inches(0.3), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), Inches(2), Inches(0.25),
          "Total COGS", size=10, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(8.84), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.25),
          "~$17.25  (5k: ~$13.50)", size=10, bold=True, align=PP_ALIGN.RIGHT)

pptx_footer(slide, 6, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 7: Hardest Problems
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Hardest Problems", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_hard = [
    ("Distinguishing shush-worthy noise from cafe ambiance",
     "The front mic picks up everything in the aimed direction -- "
     "the loud talker, but also music, dish clatter, espresso machines. "
     "The classifier must isolate sustained loud speech above the "
     "ambient baseline and hold for seconds before triggering. "
     "Too sensitive: shushes the barista. Too conservative: never fires."),
    ("Making the shush sound convincingly human",
     "If the shush sounds electronic or robotic, the product fails "
     "socially. It needs to sound like a real person -- breathy, slightly "
     "variable, naturally timed. Speaker coloration, waveguide resonances, "
     "and enclosure vibration could add artifacts. "
     "Multiple recordings with randomized selection on each firing."),
    ("Achieving useful directivity from a small speaker",
     "A 70mm puck with a 28mm speaker and short waveguide won't "
     "produce a tight beam. Target: ~4-6 dB front-to-side ratio at "
     "2-4 kHz (the \"shhh\" frequency range). Enough to make the shush "
     "louder for the target than for people at 90 degrees, but not a "
     "laser. Waveguide geometry and speaker placement are critical."),
]

for i, (title, desc) in enumerate(pptx_hard):
    y = Inches(1.15 + i * 2.4)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.15), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              f"{i+1}. {title}", size=13, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.55), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.4),
              desc, size=10, color=C_LIGHT_GRAY)

# Bottom note
pptx_rounded_rect(slide, PM, Inches(8.4), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "All three require physical prototyping. Acoustic porting, "
          "waveguide geometry, and sound recording quality can only be "
          "validated by building and testing in real environments.",
          size=11, color=C_ACCENT_ORANGE)

pptx_footer(slide, 7, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 8: Gate Result & Next
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Gate Result & Next", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

# Gate badge
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_rect(slide, PM + Inches(0.15), Inches(1.2), Inches(2.0), Inches(0.6), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.3), Inches(1.22), Inches(1.8), Inches(0.3),
          "GATE: 74 PASS", size=16, bold=True)
pptx_text(slide, PM + Inches(0.3), Inches(1.52), Inches(1.8), Inches(0.2),
          "74 pass / 6 N/A / 8 fail", size=9)
pptx_text(slide, PM + Inches(2.4), Inches(1.28), Inches(4), Inches(0.2),
          "System description complete.", size=11)
pptx_text(slide, PM + Inches(2.4), Inches(1.55), Inches(4), Inches(0.2),
          "8 items need prototype validation.", size=11, color=C_ACCENT_ORANGE, bold=True)

# Power summary
pptx_rounded_rect(slide, PM, Inches(2.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.15), SLIDE_W - 2 * PM, Inches(0.2),
          "POWER", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "84 mA active  |  2.3 mA standby  |  1200 mAh battery", size=11)
pptx_text(slide, PM + Inches(0.15), Inches(2.65), SLIDE_W - 2 * PM, Inches(0.2),
          "~3 days typical  |  USB-C charge  |  ~2.5 hr 0-100%",
          size=10, color=C_LIGHT_GRAY)

# Key specs
pptx_rounded_rect(slide, PM, Inches(3.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(3.15), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY SPECS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(3.4), SLIDE_W - 2 * PM, Inches(0.2),
          "70mm x 15mm puck  |  ~65g  |  ESP32-S3  |  BLE 5.0", size=10)
pptx_text(slide, PM + Inches(0.15), Inches(3.65), SLIDE_W - 2 * PM, Inches(0.2),
          "2 MEMS mics  |  Waveguide speaker  |  $50-70 retail",
          size=10, color=C_LIGHT_GRAY)

# Open items
pptx_rect(slide, PM, Inches(4.2), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_ORANGE)
pptx_text(slide, PM, Inches(4.3), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY OPEN ITEMS", size=11, color=C_ACCENT_ORANGE, bold=True)

pptx_open_items = [
    ("M2", "Cardioid mic porting: 10-15 dB front-to-rear rejection?"),
    ("M2", "Waveguide: >4 dB front-to-side ratio at 2-4 kHz?"),
    ("M2", "Shush naturalness through 28mm speaker in plastic enclosure"),
    ("M3", "False positive rate in 5+ real cafe environments"),
    ("M2", "ESP32-S3 BLE + dual-mic audio coexistence under load"),
    ("M2", "Sound recording session: 20+ shush variants, 5 styles"),
    ("M1", "Market positioning: $50 tool vs. $30 fun gadget"),
    ("M3", "DSP duty-cycle optimization for extended battery life"),
]

for i, (milestone, desc) in enumerate(pptx_open_items):
    y = Inches(4.6 + i * 0.47)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.38), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.07), Inches(0.5), Inches(0.2),
              milestone, size=9, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.7), y + Inches(0.07), SLIDE_W - 2 * PM - Inches(1.0), Inches(0.2),
              desc, size=9, color=C_SOFT_WHITE)

# CTA
pptx_rounded_rect(slide, PM, Inches(8.45), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.5), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "Next step: breadboard prototype on an ESP32-S3 devkit with "
          "two MEMS mics and a small speaker. Test the front/rear "
          "differential in a real cafe. Record the shush samples.",
          size=10)

pptx_footer(slide, 8, TOTAL_PAGES)

# Save PPTX
pptx_output = os.path.join(_DIR, "Shusher_Carousel.pptx")
prs.save(pptx_output)
print(f"Saved {TOTAL_PAGES}-slide carousel PPTX to {pptx_output}")
