#!/usr/bin/env python3
"""Build a LinkedIn carousel PDF and PPTX for Bubbler — Automated Soap Bubble Maker.

Format: 1080x1350 px (4:5 portrait) — optimized for mobile feed.
Uses reportlab for PDF and python-pptx for PPTX generation.
"""

import os
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor, white
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

_DIR = os.path.dirname(os.path.abspath(__file__))

# -- Page size: 4:5 ratio --
PW = 190 * mm
PH = 237.5 * mm

# -- Colors --
DARK_BG = "#1A1A2E"
CARD_BG = "#22223A"
CARD_BG_ALT = "#1E1E34"
ACCENT_ORANGE = "#FF8C00"
ACCENT_GREEN = "#4EC978"
ACCENT_RED = "#FF4545"
ACCENT_BLUE = "#009BF5"
ACCENT_PURPLE = "#A855F7"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY = "#BBBBCC"
SOFT_WHITE = "#F0F0F5"

M = 10 * mm  # standard margin


def bg(c, color=DARK_BG):
    c.setFillColor(HexColor(color))
    c.rect(0, 0, PW, PH, fill=1, stroke=0)


def accent_strip(c, color=ACCENT_ORANGE):
    c.setFillColor(HexColor(color))
    c.rect(0, PH - 2 * mm, PW, 2 * mm, fill=1, stroke=0)


def card(c, x, y, w, h, color=CARD_BG):
    """Draw a card. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    c.roundRect(x, PH - y - h, w, h, 3 * mm, fill=1, stroke=0)


def card_flat(c, x, y, w, h, color=CARD_BG):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def txt(c, x, y, text, size=14, color=WHITE_HEX, bold=False, align="left", max_w=None):
    """Draw text. y is from TOP of page."""
    c.setFillColor(HexColor(color))
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    ry = PH - y - size * 0.35
    if align == "center" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + (max_w - tw) / 2, ry, text)
    elif align == "right" and max_w:
        tw = c.stringWidth(text, font, size)
        c.drawString(x + max_w - tw, ry, text)
    else:
        c.drawString(x, ry, text)


def txt_wrap(c, x, y, text, size=11, color=WHITE_HEX, bold=False, line_h=None, max_w=None):
    """Draw wrapped text. Returns y after last line."""
    if line_h is None:
        line_h = size * 1.4
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFillColor(HexColor(color))
    c.setFont(font, size)
    if max_w is None:
        max_w = PW - 2 * M
    words = text.split()
    lines = []
    current = ""
    for w in words:
        test = current + (" " if current else "") + w
        if c.stringWidth(test, font, size) > max_w:
            if current:
                lines.append(current)
            current = w
        else:
            current = test
    if current:
        lines.append(current)
    for line in lines:
        c.drawString(x, PH - y - size * 0.35, line)
        y += line_h
    return y


def bar(c, x, y, w, h, color):
    c.setFillColor(HexColor(color))
    c.rect(x, PH - y - h, w, h, fill=1, stroke=0)


def circle_num(c, x, y, num, color):
    r = 4 * mm
    c.setFillColor(HexColor(color))
    c.circle(x + r, PH - y - r, r, fill=1, stroke=0)
    c.setFillColor(white)
    c.setFont("Helvetica-Bold", 12)
    tw = c.stringWidth(str(num), "Helvetica-Bold", 12)
    c.drawString(x + r - tw / 2, PH - y - r - 4, str(num))


def footer(c, page, total):
    c.setFillColor(HexColor(LIGHT_GRAY))
    c.setFont("Helvetica", 8)
    label = f"{page}/{total}"
    tw = c.stringWidth(label, "Helvetica", 8)
    c.drawString(PW - M - tw, 4 * mm, label)


TOTAL_PAGES = 8
output = os.path.join(_DIR, "Bubbler_Carousel.pdf")
c = canvas.Canvas(output, pagesize=(PW, PH))

# ================================================================
# PAGE 1: Title
# ================================================================
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 14 * mm, "Bubbler", size=42, color=WHITE_HEX, bold=True)

bar(c, M, 28 * mm, 35 * mm, 1 * mm, ACCENT_ORANGE)

txt_wrap(c, M, 34 * mm,
         "Big bubbles. Zero effort.",
         size=18, color=ACCENT_GREEN, max_w=PW - 2 * M)

txt_wrap(c, M, 47 * mm,
         "An automated soap bubble machine that produces giant bubbles "
         "up to 500 mm using force-curve feedback. No app, no cloud -- "
         "just press power and watch.",
         size=11, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Image
img_path = os.path.join(_DIR, "arrangement_options.png")
if os.path.exists(img_path):
    img_top_y = 68 * mm
    img_max_w = PW - 2 * M
    img_max_h = 140 * mm
    c.drawImage(ImageReader(img_path),
                M, PH - img_top_y - img_max_h,
                width=img_max_w, height=img_max_h,
                preserveAspectRatio=True, anchor='sw', mask='auto')

# Bottom bar
card_flat(c, 0, PH - 8 * mm, PW, 8 * mm, CARD_BG)
txt(c, M, PH - 5 * mm, "Product Overview  |  Concept Stage  |  2026",
    size=9, color=LIGHT_GRAY)

footer(c, 1, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 2: The Problem
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "The Problem", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

problems = [
    ("Large bubbles are hard to produce",
     "Giant soap bubbles (200-500 mm) require precise airflow control, "
     "correct dip timing, and film that survives rotation. Manual technique "
     "is inconsistent and takes practice to learn.",
     ACCENT_RED),
    ("Existing machines make only small bubbles",
     "Consumer bubble machines produce 20-50 mm bubbles with high pop rates "
     "and no adaptation. They blow hard, pop fast, and waste soap. No feedback "
     "loop means no improvement over time.",
     ACCENT_ORANGE),
    ("Gap between cheap toys and pro gear",
     "Toy machines cost $10-30 but produce tiny bubbles. Professional stage "
     "equipment costs $200+ and needs power outlets. Nothing in between serves "
     "families, performers, and outdoor events well.",
     ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(problems):
    y_top = 30 * mm + i * 55 * mm
    card(c, M, y_top, PW - 2 * M, 49 * mm, CARD_BG)
    bar(c, M, y_top, PW - 2 * M, 1.5 * mm, color)
    txt(c, M + 5 * mm, y_top + 8 * mm, title, size=14, color=color, bold=True)
    txt_wrap(c, M + 5 * mm, y_top + 20 * mm, desc,
             size=11, color=LIGHT_GRAY, max_w=PW - 2 * M - 10 * mm, line_h=14)

# Target users
card_flat(c, M, 200 * mm, PW - 2 * M, 30 * mm, CARD_BG)
txt(c, M + 4 * mm, 204 * mm, "TARGET USERS", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 213 * mm, "Families  |  Performers  |  Event organizers",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 222 * mm, "Outdoor parties  |  Buskers  |  Kids' entertainment",
    size=12, color=WHITE_HEX, bold=True)

footer(c, 2, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 3: How It Works
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "How It Works", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

steps = [
    ("Fill", "Fill vat with soap solution",
     "Pour bubble solution into the built-in vat. Standard dish soap "
     "mix works. The wand loop sits submerged, ready to dip.",
     ACCENT_BLUE),
    ("Press", "Press power -- wand dips into vat",
     "Single button press starts the cycle. The motor arm dips the "
     "wand loop into the soap solution, coating it with a thin film.",
     ACCENT_GREEN),
    ("Inflate", "Arm rotates up, fan gently inflates film",
     "The arm rotates 175 degrees upward. A DC fan blows a controlled "
     "ramp of air through the soap film, inflating it into a large bubble.",
     ACCENT_PURPLE),
    ("Optimize", "Force sensing optimizes each cycle",
     "A strain gauge on the wand arm measures 50-200 mN during inflation. "
     "Firmware classifies each outcome and auto-adjusts fan speed, dip "
     "duration, and blow ramp. Converges in 5-10 cycles.",
     ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(steps):
    y_top = 28 * mm + i * 50 * mm
    card(c, M, y_top, PW - 2 * M, 44 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, color)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=18, color=color, bold=True)
    txt(c, M + 16 * mm, y_top + 16 * mm, subtitle, size=11, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 26 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom note
card_flat(c, M, 232 * mm, PW - 2 * M, 4 * mm, CARD_BG)

footer(c, 3, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 4: Architecture
# ================================================================
bg(c)
accent_strip(c, ACCENT_BLUE)

txt(c, M, 12 * mm, "Architecture", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_BLUE)

# Signal chain
card(c, M, 28 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 32 * mm, "Sensing Signal Chain", size=11, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 41 * mm, "Strain gauge --> HX711 ADC --> STM32 MCU",
    size=12, color=WHITE_HEX, bold=True)
txt(c, M + 4 * mm, 49 * mm, "50-200 mN force       24-bit 10-80 Hz       hill-climbing optimizer",
    size=8, color=LIGHT_GRAY)

# Actuation chain
card(c, M, 60 * mm, PW - 2 * M, 22 * mm, CARD_BG)
txt(c, M + 4 * mm, 64 * mm, "Actuation Chain", size=11, color=ACCENT_GREEN, bold=True)
txt(c, M + 4 * mm, 72 * mm, "STM32 --> Motor H-bridge (dip/rotate) + Fan PWM (inflate)",
    size=10, color=WHITE_HEX, bold=True)

# Subsystems
bar(c, M, 88 * mm, PW - 2 * M, 1 * mm, ACCENT_BLUE)
txt(c, M, 92 * mm, "KEY COMPONENTS", size=11, color=ACCENT_BLUE, bold=True)

subsystems = [
    ("MCU (STM32)", "Bare-metal firmware, hill-climbing on 5 parameters, no app/cloud", ACCENT_GREEN),
    ("Strain Gauge + HX711", "Force-curve feedback at 10-80 Hz, cycle outcome classification", ACCENT_PURPLE),
    ("DC Motor + H-bridge", "Wand dip and 175-degree arm rotation, bidirectional control", ACCENT_GREEN),
    ("DC Fan + MOSFET", "PWM-controlled airflow ramp for gentle bubble inflation", ACCENT_BLUE),
    ("Power (4xAA batteries)", "~7 hr runtime, simple replacement, no charging needed", ACCENT_ORANGE),
    ("Enclosure (IPX4)", "215x206 mm footprint, ~250 mm height, splash-resistant", ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(subsystems):
    y_top = 98 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 4 * mm, name, size=11, color=WHITE_HEX, bold=True)
    txt(c, M + 4 * mm, y_top + 11 * mm, desc, size=8, color=LIGHT_GRAY)
    bar(c, M, y_top, 1.5 * mm, 14 * mm, color)

# Arrangement diagram
arr_path = os.path.join(_DIR, "arrangement_options.png")
if os.path.exists(arr_path):
    card(c, M, 202 * mm, PW - 2 * M, 30 * mm, CARD_BG_ALT)
    txt(c, M + 4 * mm, 205 * mm, "Component Arrangement", size=10, color=ACCENT_PURPLE, bold=True)
    c.drawImage(ImageReader(arr_path),
                M + 4 * mm, PH - 230 * mm,
                width=PW - 2 * M - 8 * mm, height=20 * mm,
                preserveAspectRatio=True, anchor='sw', mask='auto')

footer(c, 4, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 5: Key Innovation — Force-Curve Feedback
# ================================================================
bg(c)
accent_strip(c, ACCENT_PURPLE)

txt(c, M, 12 * mm, "Key Innovation", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_PURPLE)

txt_wrap(c, M, 28 * mm,
         "Force-curve feedback: the strain gauge on the wand arm turns "
         "every inflation cycle into a learning opportunity.",
         size=13, color=LIGHT_GRAY, max_w=PW - 2 * M)

# Sensing card
card(c, M, 44 * mm, PW - 2 * M, 45 * mm, CARD_BG)
bar(c, M, 44 * mm, PW - 2 * M, 1.5 * mm, ACCENT_GREEN)
txt(c, M + 5 * mm, 50 * mm, "What It Measures", size=16, color=ACCENT_GREEN, bold=True)
txt(c, M + 5 * mm, 62 * mm, "Force: 50-200 mN during inflation", size=12, color=WHITE_HEX)
txt(c, M + 5 * mm, 72 * mm, "Sample rate: 10-80 Hz via HX711 ADC", size=12, color=WHITE_HEX)
txt(c, M + 5 * mm, 82 * mm, "Resolution: sub-mN with 24-bit ADC", size=12, color=WHITE_HEX)

# Classification card
card(c, M, 96 * mm, PW - 2 * M, 55 * mm, CARD_BG)
bar(c, M, 96 * mm, PW - 2 * M, 1.5 * mm, ACCENT_PURPLE)
txt(c, M + 5 * mm, 102 * mm, "Cycle Outcome Classification", size=16, color=ACCENT_PURPLE, bold=True)

outcomes = [
    ("Success", "Bubble detaches cleanly -- force drops to baseline", ACCENT_GREEN),
    ("Pop", "Sudden force spike then zero -- film burst mid-inflation", ACCENT_RED),
    ("No film", "Near-zero force throughout -- dip failed to coat", ACCENT_ORANGE),
    ("Partial", "Force plateau then slow decay -- bubble formed but small", ACCENT_BLUE),
]

ry = 114 * mm
for label, desc, color in outcomes:
    txt(c, M + 5 * mm, ry, f"{label}:", size=10, color=color, bold=True)
    ry = txt_wrap(c, M + 25 * mm, ry, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 30 * mm, line_h=12)
    ry += 2 * mm

# Optimization card
card(c, M, 158 * mm, PW - 2 * M, 70 * mm, CARD_BG_ALT)
txt(c, M + 5 * mm, 163 * mm, "Hill-Climbing Optimizer", size=13, color=ACCENT_ORANGE, bold=True)

params = [
    ("Fan speed ramp", "PWM duty cycle profile during inflation"),
    ("Dip duration", "How long the wand stays in the soap vat"),
    ("Blow ramp rate", "How quickly airflow increases"),
    ("Rotation speed", "Arm angular velocity during lift"),
    ("Pause duration", "Wait time between dip and blow"),
]

ry = 175 * mm
for title, desc in params:
    txt(c, M + 5 * mm, ry, f"- {title}:", size=10, color=WHITE_HEX, bold=True)
    ry = txt_wrap(c, M + 8 * mm, ry + 12, desc,
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 16 * mm, line_h=12)
    ry += 2 * mm

footer(c, 5, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 6: Constraints & BOM
# ================================================================
bg(c)
accent_strip(c, ACCENT_ORANGE)

txt(c, M, 12 * mm, "Constraints & BOM", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_ORANGE)

constraints = [
    ("Bubble size", "Up to 500 mm", "Wand loop geometry + controlled airflow"),
    ("Battery life", "~7 hrs (4xAA)", "Low-power MCU, motor duty-cycled"),
    ("Footprint", "215 x 206 mm", "Compact enough for a tabletop"),
    ("Operating temp", "5-40 C", "Soap film physics limit the range"),
    ("Wind tolerance", "4 kph crosswind", "Fan ramp compensation algorithm"),
    ("Electronics", "IPX4", "Splash-resistant enclosure for outdoor use"),
]

for i, (name, value, note) in enumerate(constraints):
    y_top = 28 * mm + i * 17 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 14 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, name, size=11, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 4 * mm, y_top + 10 * mm, note, size=8, color=LIGHT_GRAY)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, value, size=11, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 8 * mm)

# BOM
bar(c, M, 133 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 137 * mm, "BOM ESTIMATE (~$12.80 total)", size=11, color=ACCENT_GREEN, bold=True)

bom = [
    ("STM32 MCU", "$1.50"),
    ("Strain gauge + HX711 ADC", "$0.80"),
    ("DC motor + H-bridge", "$1.80"),
    ("DC fan + MOSFET driver", "$1.10"),
    ("Enclosure (molded plastic)", "$3.50"),
    ("Wand + loop + shaft", "$1.50"),
    ("Misc (PCB, connectors, passives)", "$2.60"),
]

for i, (item, cost) in enumerate(bom):
    y_top = 144 * mm + i * 9 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card_flat(c, M + 2 * mm, y_top, PW - 2 * M - 4 * mm, 7.5 * mm, bg_c)
    txt(c, M + 6 * mm, y_top + 2 * mm, item, size=9, color=SOFT_WHITE)
    txt(c, M + 6 * mm, y_top + 2 * mm, cost, size=9, color=WHITE_HEX, bold=True,
        align="right", max_w=PW - 2 * M - 16 * mm)

# Total
card_flat(c, M + 2 * mm, 210 * mm, PW - 2 * M - 4 * mm, 7 * mm, ACCENT_ORANGE)
txt(c, M + 6 * mm, 212 * mm, "Total BOM", size=10, color=WHITE_HEX, bold=True)
txt(c, M + 6 * mm, 212 * mm, "~$12.80  |  Target retail: sub-$50", size=10, color=WHITE_HEX, bold=True,
    align="right", max_w=PW - 2 * M - 16 * mm)

# Retail note
txt_wrap(c, M, 222 * mm,
         "4xAA batteries not included in BOM. No app, no cloud -- "
         "keeps ongoing costs at zero.",
         size=9, color=LIGHT_GRAY, max_w=PW - 2 * M)

footer(c, 6, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 7: Hardest Problems
# ================================================================
bg(c)
accent_strip(c, ACCENT_RED)

txt(c, M, 12 * mm, "Hardest Problems", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_RED)

hard_problems = [
    ("Film survival during rotation",
     "The soap film must survive a 175-degree arm rotation from the vat "
     "to the blow position without breaking. Film thickness, rotation "
     "speed, and acceleration profile all matter. Too fast and the film "
     "tears from inertia; too slow and it drains and thins. Requires "
     "empirical tuning of the motor ramp curve."),
    ("Force-curve interpretation",
     "Classifying cycle outcomes (success, pop, no-film, partial) from "
     "noisy strain gauge signals at 10-80 Hz. The HX711 output includes "
     "mechanical vibration, motor coupling, and wind noise. Signal "
     "conditioning and threshold-based classification must be reliable "
     "enough for the hill-climbing optimizer to converge."),
    ("Wind compensation at 4 kph",
     "Outdoor use means crosswind. A 4 kph breeze changes the effective "
     "airflow through the soap film, shifting optimal fan speed and blow "
     "duration. The optimizer must detect wind-induced pop patterns and "
     "compensate within a few cycles. No wind sensor -- inferred from "
     "force-curve anomalies only."),
]

for i, (title, desc) in enumerate(hard_problems):
    y_top = 30 * mm + i * 62 * mm
    card(c, M, y_top, PW - 2 * M, 56 * mm, CARD_BG)
    circle_num(c, M + 4 * mm, y_top + 4 * mm, i + 1, ACCENT_RED)
    txt(c, M + 16 * mm, y_top + 6 * mm, title, size=13, color=WHITE_HEX, bold=True)
    txt_wrap(c, M + 6 * mm, y_top + 18 * mm, desc,
             size=10, color=LIGHT_GRAY, max_w=PW - 2 * M - 12 * mm, line_h=13)

# Bottom
card_flat(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "All three require physical prototyping. Soap film behavior "
         "cannot be fully simulated -- build, measure, iterate.",
         size=11, color=ACCENT_ORANGE, max_w=PW - 2 * M - 8 * mm, line_h=14)

footer(c, 7, TOTAL_PAGES)
c.showPage()

# ================================================================
# PAGE 8: Gate Result & Next
# ================================================================
bg(c)
accent_strip(c, ACCENT_GREEN)

txt(c, M, 12 * mm, "Gate Result & Next", size=28, color=WHITE_HEX, bold=True)
bar(c, M, 22 * mm, 30 * mm, 0.8 * mm, ACCENT_GREEN)

# Gate badge
card(c, M, 28 * mm, PW - 2 * M, 22 * mm, CARD_BG)
card(c, M + 4 * mm, 31 * mm, 50 * mm, 16 * mm, ACCENT_GREEN)
txt(c, M + 8 * mm, 35 * mm, "GATE: PASS", size=16, color=WHITE_HEX, bold=True)
txt(c, M + 8 * mm, 43 * mm, "62 pass / 23 N/A / 3 minor", size=9, color=WHITE_HEX)
txt(c, M + 60 * mm, 36 * mm, "System description complete.", size=11, color=WHITE_HEX)
txt(c, M + 60 * mm, 44 * mm, "Ready to proceed to PRD.", size=11, color=ACCENT_GREEN, bold=True)

# Key specs
card(c, M, 54 * mm, PW - 2 * M, 28 * mm, CARD_BG)
txt(c, M + 4 * mm, 58 * mm, "KEY SPECS", size=10, color=ACCENT_BLUE, bold=True)
txt(c, M + 4 * mm, 66 * mm, "Bubbles up to 500 mm  |  4xAA (~7 hr)  |  215x206 mm footprint",
    size=10, color=WHITE_HEX)
txt(c, M + 4 * mm, 74 * mm, "~250 mm height  |  IPX4  |  5-40 C  |  BOM ~$12.80  |  Sub-$50 retail",
    size=10, color=LIGHT_GRAY)

# Architecture summary
card(c, M, 86 * mm, PW - 2 * M, 18 * mm, CARD_BG_ALT)
txt(c, M + 4 * mm, 90 * mm, "ARCHITECTURE", size=10, color=ACCENT_PURPLE, bold=True)
txt(c, M + 4 * mm, 98 * mm, "Strain gauge + HX711 + STM32  |  Bare-metal FW  |  No app, no cloud",
    size=10, color=WHITE_HEX)

# Minor gaps
bar(c, M, 110 * mm, PW - 2 * M, 1 * mm, ACCENT_ORANGE)
txt(c, M, 114 * mm, "3 MINOR GAPS", size=11, color=ACCENT_ORANGE, bold=True)

gaps = [
    ("Gap 1", "FW versioning scheme not yet defined"),
    ("Gap 2", "Decision consequences formatting incomplete"),
    ("Gap 3", "Schedule milestones need dates"),
]

for i, (label, desc) in enumerate(gaps):
    y_top = 120 * mm + i * 14 * mm
    bg_c = CARD_BG if i % 2 == 0 else CARD_BG_ALT
    card(c, M, y_top, PW - 2 * M, 11 * mm, bg_c)
    txt(c, M + 4 * mm, y_top + 3.5 * mm, label, size=9, color=ACCENT_ORANGE, bold=True)
    txt(c, M + 22 * mm, y_top + 3.5 * mm, desc, size=9, color=SOFT_WHITE)

# What's next
bar(c, M, 166 * mm, PW - 2 * M, 1 * mm, ACCENT_GREEN)
txt(c, M, 170 * mm, "WHAT'S NEXT", size=11, color=ACCENT_GREEN, bold=True)

next_items = [
    "Build a functional prototype with off-the-shelf motor, fan, and HX711 breakout",
    "Validate soap film survival during arm rotation (175 deg)",
    "Test force-curve classification accuracy across soap formulations",
    "Measure wind compensation convergence at 4 kph crosswind",
    "Confirm 4xAA battery life target with real duty cycles",
]

ry = 178 * mm
for item in next_items:
    ry = txt_wrap(c, M + 4 * mm, ry, f"- {item}",
                  size=9, color=LIGHT_GRAY, max_w=PW - 2 * M - 8 * mm, line_h=12)
    ry += 1 * mm

# CTA
card(c, M, 218 * mm, PW - 2 * M, 16 * mm, CARD_BG)
txt_wrap(c, M + 4 * mm, 222 * mm,
         "Next step: build the mechanical prototype. Validate film survival, "
         "force sensing, and the optimization loop before committing to PCB.",
         size=10, color=WHITE_HEX, max_w=PW - 2 * M - 8 * mm, line_h=13)

footer(c, 8, TOTAL_PAGES)
c.showPage()

# ================================================================
c.save()
print(f"Saved {TOTAL_PAGES}-page carousel PDF to {output}")
print(f"Page size: {PW/mm:.0f} x {PH/mm:.0f} mm (4:5 ratio)")

# ================================================================
# PPTX GENERATION
# ================================================================
# 4:5 portrait slide: 7.5 x 9.375 inches (same ratio as PDF)
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(9.375)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# -- PPTX helper colors --
C_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
C_CARD_BG = RGBColor(0x22, 0x22, 0x3A)
C_CARD_BG_ALT = RGBColor(0x1E, 0x1E, 0x34)
C_ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
C_ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x78)
C_ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
C_ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
C_ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

PM = Inches(0.4)  # margin


def pptx_bg(slide, color=C_DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def pptx_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def pptx_text(slide, left, top, width, height, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    return txBox


def pptx_add_para(tf, text, size=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Helvetica"
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p


def pptx_footer(slide, page, total):
    pptx_text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.3),
              Inches(0.6), Inches(0.25), f"{page}/{total}",
              size=8, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ================================================================
# PPTX PAGE 1: Title
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.5), SLIDE_W - 2 * PM, Inches(0.7),
          "Bubbler", size=42, bold=True)

pptx_rect(slide, PM, Inches(1.1), Inches(1.4), Inches(0.04), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(1.3), SLIDE_W - 2 * PM, Inches(0.4),
          "Big bubbles. Zero effort.", size=18, color=C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(1.8), SLIDE_W - 2 * PM, Inches(0.6),
          "An automated soap bubble machine that produces giant bubbles "
          "up to 500 mm using force-curve feedback. No app, no cloud -- "
          "just press power and watch.",
          size=11, color=C_LIGHT_GRAY)

# Image
img_path = os.path.join(_DIR, "arrangement_options.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, PM, Inches(2.7),
                             width=SLIDE_W - 2 * PM, height=Inches(5.5))

pptx_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), C_CARD_BG)
pptx_text(slide, PM, SLIDE_H - Inches(0.3), SLIDE_W - 2 * PM, Inches(0.2),
          "Product Overview  |  Concept Stage  |  2026", size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 1, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 2: The Problem
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "The Problem", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_problems = [
    ("Large bubbles are hard to produce",
     "Giant soap bubbles (200-500 mm) require precise airflow control, "
     "correct dip timing, and film that survives rotation. Manual technique "
     "is inconsistent and takes practice to learn.",
     C_ACCENT_RED),
    ("Existing machines make only small bubbles",
     "Consumer bubble machines produce 20-50 mm bubbles with high pop rates "
     "and no adaptation. They blow hard, pop fast, and waste soap. No feedback "
     "loop means no improvement over time.",
     C_ACCENT_ORANGE),
    ("Gap between cheap toys and pro gear",
     "Toy machines cost $10-30 but produce tiny bubbles. Professional stage "
     "equipment costs $200+ and needs power outlets. Nothing in between serves "
     "families, performers, and outdoor events well.",
     C_ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(pptx_problems):
    y = Inches(1.15 + i * 2.15)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.9), C_CARD_BG)
    pptx_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.06), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.3),
              title, size=14, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.65), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
              desc, size=11, color=C_LIGHT_GRAY)

# Target users
pptx_rounded_rect(slide, PM, Inches(7.7), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
          "TARGET USERS", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.1), SLIDE_W - 2 * PM, Inches(0.25),
          "Families  |  Performers  |  Event organizers", size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Outdoor parties  |  Buskers  |  Kids' entertainment", size=12, bold=True)

pptx_footer(slide, 2, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 3: How It Works
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "How It Works", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

pptx_steps = [
    ("1. Fill", "Fill vat with soap solution",
     "Pour bubble solution into the built-in vat. Standard dish soap "
     "mix works. The wand loop sits submerged, ready to dip.",
     C_ACCENT_BLUE),
    ("2. Press", "Press power -- wand dips into vat",
     "Single button press starts the cycle. The motor arm dips the "
     "wand loop into the soap solution, coating it with a thin film.",
     C_ACCENT_GREEN),
    ("3. Inflate", "Arm rotates up, fan gently inflates film",
     "The arm rotates 175 degrees upward. A DC fan blows a controlled "
     "ramp of air through the soap film, inflating it into a large bubble.",
     C_ACCENT_PURPLE),
    ("4. Optimize", "Force sensing optimizes each cycle",
     "A strain gauge on the wand arm measures 50-200 mN during inflation. "
     "Firmware classifies each outcome and auto-adjusts fan speed, dip "
     "duration, and blow ramp. Converges in 5-10 cycles.",
     C_ACCENT_ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(pptx_steps):
    y = Inches(1.1 + i * 1.95)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.1), Inches(3), Inches(0.35),
              title, size=18, color=color, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.45), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              subtitle, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.8), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.8),
              desc, size=10, color=C_LIGHT_GRAY)

pptx_footer(slide, 3, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 4: Architecture
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_BLUE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Architecture", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_BLUE)

# Signal chain
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(1.1), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(1.15), SLIDE_W - 2 * PM, Inches(0.2),
          "Sensing Signal Chain", size=11, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.45), SLIDE_W - 2 * PM, Inches(0.25),
          "Strain gauge --> HX711 ADC --> STM32 MCU",
          size=12, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(1.8), SLIDE_W - 2 * PM, Inches(0.2),
          "50-200 mN force       24-bit 10-80 Hz       hill-climbing optimizer",
          size=8, color=C_LIGHT_GRAY)

# Actuation chain
pptx_rounded_rect(slide, PM, Inches(2.35), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Actuation Chain", size=11, color=C_ACCENT_GREEN, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.7), SLIDE_W - 2 * PM, Inches(0.25),
          "STM32 --> Motor H-bridge (dip/rotate) + Fan PWM (inflate)",
          size=10, bold=True)

pptx_rect(slide, PM, Inches(3.4), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_BLUE)
pptx_text(slide, PM, Inches(3.5), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY COMPONENTS", size=11, color=C_ACCENT_BLUE, bold=True)

pptx_subsystems = [
    ("MCU (STM32)", "Bare-metal firmware, hill-climbing on 5 parameters, no app/cloud", C_ACCENT_GREEN),
    ("Strain Gauge + HX711", "Force-curve feedback at 10-80 Hz, cycle outcome classification", C_ACCENT_PURPLE),
    ("DC Motor + H-bridge", "Wand dip and 175-degree arm rotation, bidirectional control", C_ACCENT_GREEN),
    ("DC Fan + MOSFET", "PWM-controlled airflow ramp for gentle bubble inflation", C_ACCENT_BLUE),
    ("Power (4xAA batteries)", "~7 hr runtime, simple replacement, no charging needed", C_ACCENT_ORANGE),
    ("Enclosure (IPX4)", "215x206 mm footprint, ~250 mm height, splash-resistant", C_ACCENT_ORANGE),
]

for i, (name, desc, color) in enumerate(pptx_subsystems):
    y = Inches(3.8 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_rect(slide, PM, y, Inches(0.06), Inches(0.55), color)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), SLIDE_W - 2 * PM, Inches(0.2),
              name, size=11, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.28), SLIDE_W - 2 * PM, Inches(0.2),
              desc, size=8, color=C_LIGHT_GRAY)

# Arrangement diagram
arr_path = os.path.join(_DIR, "arrangement_options.png")
if os.path.exists(arr_path):
    pptx_rounded_rect(slide, PM, Inches(7.75), SLIDE_W - 2 * PM, Inches(1.2), C_CARD_BG_ALT)
    pptx_text(slide, PM + Inches(0.15), Inches(7.8), SLIDE_W - 2 * PM, Inches(0.2),
              "Component Arrangement", size=10, color=C_ACCENT_PURPLE, bold=True)
    slide.shapes.add_picture(arr_path, PM + Inches(0.15), Inches(8.1),
                             width=SLIDE_W - 2 * PM - Inches(0.3), height=Inches(0.75))

pptx_footer(slide, 4, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 5: Key Innovation
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Key Innovation", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_PURPLE)

pptx_text(slide, PM, Inches(1.05), SLIDE_W - 2 * PM, Inches(0.5),
          "Force-curve feedback: the strain gauge on the wand arm turns "
          "every inflation cycle into a learning opportunity.",
          size=13, color=C_LIGHT_GRAY)

# Sensing card
pptx_rounded_rect(slide, PM, Inches(1.7), SLIDE_W - 2 * PM, Inches(1.7), C_CARD_BG)
pptx_rect(slide, PM, Inches(1.7), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.2), Inches(1.85), SLIDE_W - 2 * PM, Inches(0.3),
          "What It Measures", size=16, color=C_ACCENT_GREEN, bold=True)
txBox = pptx_text(slide, PM + Inches(0.2), Inches(2.3), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.0),
                  "Force: 50-200 mN during inflation", size=12)
tf = txBox.text_frame
pptx_add_para(tf, "Sample rate: 10-80 Hz via HX711 ADC", size=12, space_before=6)
pptx_add_para(tf, "Resolution: sub-mN with 24-bit ADC", size=12, space_before=6)

# Classification card
pptx_rounded_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(2.2), C_CARD_BG)
pptx_rect(slide, PM, Inches(3.6), SLIDE_W - 2 * PM, Inches(0.06), C_ACCENT_PURPLE)
pptx_text(slide, PM + Inches(0.2), Inches(3.75), SLIDE_W - 2 * PM, Inches(0.3),
          "Cycle Outcome Classification", size=16, color=C_ACCENT_PURPLE, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(4.2), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.5),
                  "Success: Bubble detaches cleanly -- force drops to baseline", size=10, color=C_ACCENT_GREEN)
tf = txBox.text_frame
pptx_add_para(tf, "Pop: Sudden force spike then zero -- film burst mid-inflation",
              size=10, color=C_ACCENT_RED, space_before=6)
pptx_add_para(tf, "No film: Near-zero force throughout -- dip failed to coat",
              size=10, color=C_ACCENT_ORANGE, space_before=6)
pptx_add_para(tf, "Partial: Force plateau then slow decay -- bubble formed but small",
              size=10, color=C_ACCENT_BLUE, space_before=6)

# Optimization card
pptx_rounded_rect(slide, PM, Inches(6.1), SLIDE_W - 2 * PM, Inches(2.8), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.2), Inches(6.2), SLIDE_W - 2 * PM, Inches(0.25),
          "Hill-Climbing Optimizer (5 parameters)", size=13, color=C_ACCENT_ORANGE, bold=True)

txBox = pptx_text(slide, PM + Inches(0.2), Inches(6.6), SLIDE_W - 2 * PM - Inches(0.4), Inches(2.0),
                  "Fan speed ramp: PWM duty cycle profile during inflation", size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
pptx_add_para(tf, "Dip duration: How long the wand stays in the soap vat",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Blow ramp rate: How quickly airflow increases",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Rotation speed: Arm angular velocity during lift",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "Pause duration: Wait time between dip and blow",
              size=9, color=C_LIGHT_GRAY, space_before=6)
pptx_add_para(tf, "", size=6, space_before=6)
pptx_add_para(tf, "Converges in 5-10 cycles to local optimum for current conditions.",
              size=10, color=C_ACCENT_GREEN, bold=True, space_before=4)

pptx_footer(slide, 5, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 6: Constraints & BOM
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_ORANGE)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Constraints & BOM", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_ORANGE)

pptx_constraints = [
    ("Bubble size", "Up to 500 mm", "Wand loop geometry + controlled airflow"),
    ("Battery life", "~7 hrs (4xAA)", "Low-power MCU, motor duty-cycled"),
    ("Footprint", "215 x 206 mm", "Compact enough for a tabletop"),
    ("Operating temp", "5-40 C", "Soap film physics limit the range"),
    ("Wind tolerance", "4 kph crosswind", "Fan ramp compensation algorithm"),
    ("Electronics", "IPX4", "Splash-resistant enclosure for outdoor use"),
]

for i, (name, value, note) in enumerate(pptx_constraints):
    y = Inches(1.1 + i * 0.65)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.55), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), Inches(3), Inches(0.2),
              name, size=11, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.3), Inches(5), Inches(0.2),
              note, size=8, color=C_LIGHT_GRAY)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.05), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.2),
              value, size=11, bold=True, align=PP_ALIGN.RIGHT)

# BOM
pptx_rect(slide, PM, Inches(5.15), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(5.25), SLIDE_W - 2 * PM, Inches(0.2),
          "BOM ESTIMATE (~$12.80 total)", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_bom = [
    ("STM32 MCU", "$1.50"),
    ("Strain gauge + HX711 ADC", "$0.80"),
    ("DC motor + H-bridge", "$1.80"),
    ("DC fan + MOSFET driver", "$1.10"),
    ("Enclosure (molded plastic)", "$3.50"),
    ("Wand + loop + shaft", "$1.50"),
    ("Misc (PCB, connectors, passives)", "$2.60"),
]

for i, (item, cost) in enumerate(pptx_bom):
    y = Inches(5.55 + i * 0.35)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rect(slide, PM + Inches(0.08), y, SLIDE_W - 2 * PM - Inches(0.16), Inches(0.28), bg_c)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), Inches(4), Inches(0.2),
              item, size=9, color=C_SOFT_WHITE)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.03), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.2),
              cost, size=9, bold=True, align=PP_ALIGN.RIGHT)

# Total
pptx_rect(slide, PM + Inches(0.08), Inches(8.05), SLIDE_W - 2 * PM - Inches(0.16), Inches(0.3), C_ACCENT_ORANGE)
pptx_text(slide, PM + Inches(0.2), Inches(8.07), Inches(2), Inches(0.25),
          "Total BOM", size=10, bold=True)
pptx_text(slide, PM + Inches(0.2), Inches(8.07), SLIDE_W - 2 * PM - Inches(0.6), Inches(0.25),
          "~$12.80  |  Target retail: sub-$50", size=10, bold=True, align=PP_ALIGN.RIGHT)

pptx_text(slide, PM, Inches(8.5), SLIDE_W - 2 * PM, Inches(0.4),
          "4xAA batteries not included in BOM. No app, no cloud -- keeps ongoing costs at zero.",
          size=9, color=C_LIGHT_GRAY)

pptx_footer(slide, 6, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 7: Hardest Problems
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_RED)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Hardest Problems", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_RED)

pptx_hard = [
    ("Film survival during rotation",
     "The soap film must survive a 175-degree arm rotation from the vat "
     "to the blow position without breaking. Film thickness, rotation "
     "speed, and acceleration profile all matter. Too fast and the film "
     "tears from inertia; too slow and it drains and thins. Requires "
     "empirical tuning of the motor ramp curve."),
    ("Force-curve interpretation",
     "Classifying cycle outcomes (success, pop, no-film, partial) from "
     "noisy strain gauge signals at 10-80 Hz. The HX711 output includes "
     "mechanical vibration, motor coupling, and wind noise. Signal "
     "conditioning and threshold-based classification must be reliable "
     "enough for the hill-climbing optimizer to converge."),
    ("Wind compensation at 4 kph",
     "Outdoor use means crosswind. A 4 kph breeze changes the effective "
     "airflow through the soap film, shifting optimal fan speed and blow "
     "duration. The optimizer must detect wind-induced pop patterns and "
     "compensate within a few cycles. No wind sensor -- inferred from "
     "force-curve anomalies only."),
]

for i, (title, desc) in enumerate(pptx_hard):
    y = Inches(1.15 + i * 2.4)
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(2.15), C_CARD_BG)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.15), SLIDE_W - 2 * PM - Inches(0.4), Inches(0.25),
              f"{i+1}. {title}", size=13, bold=True)
    pptx_text(slide, PM + Inches(0.2), y + Inches(0.55), SLIDE_W - 2 * PM - Inches(0.4), Inches(1.4),
              desc, size=10, color=C_LIGHT_GRAY)

# Bottom note
pptx_rounded_rect(slide, PM, Inches(8.4), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.45), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "All three require physical prototyping. Soap film behavior "
          "cannot be fully simulated -- build, measure, iterate.",
          size=11, color=C_ACCENT_ORANGE)

pptx_footer(slide, 7, TOTAL_PAGES)

# ================================================================
# PPTX PAGE 8: Gate Result & Next
# ================================================================
slide = prs.slides.add_slide(blank_layout)
pptx_bg(slide)
pptx_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT_GREEN)

pptx_text(slide, PM, Inches(0.35), SLIDE_W - 2 * PM, Inches(0.5),
          "Gate Result & Next", size=28, bold=True)
pptx_rect(slide, PM, Inches(0.85), Inches(1.2), Inches(0.03), C_ACCENT_GREEN)

# Gate badge
pptx_rounded_rect(slide, PM, Inches(1.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_rect(slide, PM + Inches(0.15), Inches(1.2), Inches(2.0), Inches(0.6), C_ACCENT_GREEN)
pptx_text(slide, PM + Inches(0.3), Inches(1.22), Inches(1.8), Inches(0.3),
          "GATE: PASS", size=16, bold=True)
pptx_text(slide, PM + Inches(0.3), Inches(1.52), Inches(1.8), Inches(0.2),
          "62 pass / 23 N/A / 3 minor", size=9)
pptx_text(slide, PM + Inches(2.4), Inches(1.28), Inches(4), Inches(0.2),
          "System description complete.", size=11)
pptx_text(slide, PM + Inches(2.4), Inches(1.55), Inches(4), Inches(0.2),
          "Ready to proceed to PRD.", size=11, color=C_ACCENT_GREEN, bold=True)

# Key specs
pptx_rounded_rect(slide, PM, Inches(2.1), SLIDE_W - 2 * PM, Inches(0.85), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(2.15), SLIDE_W - 2 * PM, Inches(0.2),
          "KEY SPECS", size=10, color=C_ACCENT_BLUE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(2.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Bubbles up to 500 mm  |  4xAA (~7 hr)  |  215x206 mm footprint", size=10)
pptx_text(slide, PM + Inches(0.15), Inches(2.65), SLIDE_W - 2 * PM, Inches(0.2),
          "~250 mm height  |  IPX4  |  5-40 C  |  BOM ~$12.80  |  Sub-$50 retail",
          size=10, color=C_LIGHT_GRAY)

# Architecture summary
pptx_rounded_rect(slide, PM, Inches(3.1), SLIDE_W - 2 * PM, Inches(0.7), C_CARD_BG_ALT)
pptx_text(slide, PM + Inches(0.15), Inches(3.15), SLIDE_W - 2 * PM, Inches(0.2),
          "ARCHITECTURE", size=10, color=C_ACCENT_PURPLE, bold=True)
pptx_text(slide, PM + Inches(0.15), Inches(3.4), SLIDE_W - 2 * PM, Inches(0.2),
          "Strain gauge + HX711 + STM32  |  Bare-metal FW  |  No app, no cloud", size=10)

# Minor gaps
pptx_rect(slide, PM, Inches(4.05), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_ORANGE)
pptx_text(slide, PM, Inches(4.15), SLIDE_W - 2 * PM, Inches(0.2),
          "3 MINOR GAPS", size=11, color=C_ACCENT_ORANGE, bold=True)

pptx_gaps = [
    ("Gap 1", "FW versioning scheme not yet defined"),
    ("Gap 2", "Decision consequences formatting incomplete"),
    ("Gap 3", "Schedule milestones need dates"),
]

for i, (label, desc) in enumerate(pptx_gaps):
    y = Inches(4.5 + i * 0.47)
    bg_c = C_CARD_BG if i % 2 == 0 else C_CARD_BG_ALT
    pptx_rounded_rect(slide, PM, y, SLIDE_W - 2 * PM, Inches(0.38), bg_c)
    pptx_text(slide, PM + Inches(0.15), y + Inches(0.07), Inches(0.6), Inches(0.2),
              label, size=9, color=C_ACCENT_ORANGE, bold=True)
    pptx_text(slide, PM + Inches(0.7), y + Inches(0.07), SLIDE_W - 2 * PM - Inches(1.0), Inches(0.2),
              desc, size=9, color=C_SOFT_WHITE)

# What's next
pptx_rect(slide, PM, Inches(6.0), SLIDE_W - 2 * PM, Inches(0.03), C_ACCENT_GREEN)
pptx_text(slide, PM, Inches(6.1), SLIDE_W - 2 * PM, Inches(0.2),
          "WHAT'S NEXT", size=11, color=C_ACCENT_GREEN, bold=True)

pptx_next = [
    "Build functional prototype with off-the-shelf motor, fan, and HX711 breakout",
    "Validate soap film survival during arm rotation (175 deg)",
    "Test force-curve classification accuracy across soap formulations",
    "Measure wind compensation convergence at 4 kph crosswind",
    "Confirm 4xAA battery life target with real duty cycles",
]

txBox = pptx_text(slide, PM + Inches(0.15), Inches(6.4), SLIDE_W - 2 * PM - Inches(0.3), Inches(1.5),
                  f"- {pptx_next[0]}", size=9, color=C_LIGHT_GRAY)
tf = txBox.text_frame
for item in pptx_next[1:]:
    pptx_add_para(tf, f"- {item}", size=9, color=C_LIGHT_GRAY, space_before=4)

# CTA
pptx_rounded_rect(slide, PM, Inches(8.45), SLIDE_W - 2 * PM, Inches(0.6), C_CARD_BG)
pptx_text(slide, PM + Inches(0.15), Inches(8.5), SLIDE_W - 2 * PM - Inches(0.3), Inches(0.5),
          "Next step: build the mechanical prototype. Validate film survival, "
          "force sensing, and the optimization loop before committing to PCB.",
          size=10)

pptx_footer(slide, 8, TOTAL_PAGES)

# Save PPTX
pptx_output = os.path.join(_DIR, "Bubbler_Carousel.pptx")
prs.save(pptx_output)
print(f"Saved {TOTAL_PAGES}-slide carousel PPTX to {pptx_output}")
