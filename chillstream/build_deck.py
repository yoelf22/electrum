#!/usr/bin/env python3
"""Build ChillStream executive pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# -- Theme colors --
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x9B, 0xF5)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED_ACCENT = RGBColor(0xFF, 0x45, 0x45)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if opacity is not None:
        # Set transparency via alpha
        shape.fill.fore_color.brightness = 0
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, size=16, color=WHITE, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.8), color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide)

# Accent stripe at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

# Product name
tb = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5))
set_text(tb.text_frame, "ChillStream", size=60, color=WHITE, bold=True)
add_paragraph(tb.text_frame, "Phase-Change Thermal Battery Water Dispenser", size=28, color=ACCENT_CYAN, bold=False, space_before=Pt(12))

# Tagline
tb2 = add_text_box(slide, Inches(1), Inches(4.0), Inches(9), Inches(1))
set_text(tb2.text_frame, "Cold water for the 50th person in line \u2014 not just the first.",
         size=22, color=LIGHT_GRAY, bold=False)

# Bottom bar with context
add_shape(slide, Inches(0), Inches(6.5), W, Inches(0.005), RGBColor(0x33, 0x33, 0x55))
tb3 = add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.6))
set_text(tb3.text_frame, "Executive Overview  \u2022  Product Investment Proposal  \u2022  2026",
         size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

# Title
tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "The Problem", size=36, color=WHITE, bold=True)

# Left column - the story
add_accent_bar(slide, Inches(0.8), Inches(1.6), height=Inches(4.2), color=RED_ACCENT)

tb = add_text_box(slide, Inches(1.2), Inches(1.6), Inches(5.5), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Water coolers fail at the moment of peak demand", size=22, color=WHITE, bold=True)
add_bullet(tf, "School lunch bell rings \u2014 60 students, 5 minutes", size=17, color=LIGHT_GRAY)
add_bullet(tf, "Office break time \u2014 30 people cluster at the cooler", size=17, color=LIGHT_GRAY)
add_bullet(tf, "Traditional cooler: 2-5L of pre-chilled water", size=17, color=LIGHT_GRAY)
add_bullet(tf, "After 10-15 pours, reservoir is depleted", size=17, color=LIGHT_GRAY)
add_bullet(tf, "Incoming mains water at 25\u00b0C mixes with remaining cold", size=17, color=LIGHT_GRAY)
add_bullet(tf, "Recovery time: 20-40 minutes", size=17, color=LIGHT_GRAY)

add_paragraph(tf, "", size=10, color=WHITE)
add_paragraph(tf, "The first users get cold water.", size=20, color=WHITE, bold=True)
add_paragraph(tf, "Everyone else gets lukewarm.", size=20, color=RED_ACCENT, bold=True)

# Right column - stats boxes
box_left = Inches(7.5)
box_w = Inches(4.8)
box_h = Inches(1.15)

for i, (stat, desc) in enumerate([
    ("20-40 min", "Recovery time after peak demand depletes the reservoir"),
    ("85%+", "Of pours during rush hour served above target temperature"),
    ("$0", "Revenue from solving this \u2014 no product addresses peak demand in countertop form factor"),
]):
    y = Inches(1.6) + Inches(1.4) * i
    box = add_shape(slide, box_left, y, box_w, box_h, RGBColor(0x22, 0x22, 0x3A))
    tb = add_text_box(slide, box_left + Inches(0.3), y + Inches(0.15), Inches(4.2), Inches(0.5))
    set_text(tb.text_frame, stat, size=30, color=ORANGE, bold=True)
    tb2 = add_text_box(slide, box_left + Inches(0.3), y + Inches(0.65), Inches(4.2), Inches(0.5))
    set_text(tb2.text_frame, desc, size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "The Solution: Phase-Change Thermal Battery", size=36, color=WHITE, bold=True)

# Three-column layout
cols = [
    ("Store", "IDLE TIME",
     "Compressor freezes PCM coating\non submerged HX plates to -5\u00b0C.\nCold energy stored as latent heat.\nNo water movement needed.",
     ACCENT_BLUE),
    ("Deliver", "PEAK DEMAND",
     "User opens tap. Mains pressure\npushes water through PCM-coated\nplates. Water chills from 25\u00b0C to\n4\u00b0C on contact. No pump.",
     ACCENT_CYAN),
    ("Recover", "BETWEEN RUSHES",
     "Compressor recharges the PCM.\n33 min for full 5L recharge.\nPartial recharge much faster.\n50-min school period = full reset.",
     ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(cols):
    x = Inches(0.8) + Inches(4.1) * i

    # Color accent top bar
    add_shape(slide, x, Inches(1.8), Inches(3.6), Inches(0.06), color)

    # Box background
    add_shape(slide, x, Inches(1.86), Inches(3.6), Inches(4.2), RGBColor(0x22, 0x22, 0x3A))

    # Title
    tb = add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3.0), Inches(0.6))
    set_text(tb.text_frame, title, size=28, color=color, bold=True)

    # Subtitle
    tb = add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.0), Inches(0.4))
    set_text(tb.text_frame, subtitle, size=12, color=LIGHT_GRAY, bold=True)

    # Description
    tb = add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(3.0), Inches(2.5))
    set_text(tb.text_frame, desc, size=16, color=SOFT_WHITE)

# Bottom callout
tb = add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.7))
set_text(tb.text_frame,
         "Latent heat stores 5-10x more energy per kg than sensible heat \u2014 enabling countertop form factor with 50-pour capacity.",
         size=16, color=ACCENT_CYAN, bold=False)

# ============================================================
# SLIDE 4: How It Works - Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Product Architecture", size=36, color=WHITE, bold=True)

# Three component boxes
components = [
    ("Insulated Tank", "Simple SS vessel\nHolds water at rest\nOpen top, sealed by cover\nServiceable \u2014 lift cover to clean",
     "No pressure rating needed\nNo embedded channels\nFood-grade stainless steel"),
    ("Submerged Plate HX", "Brazed SS plates with\nPCM coating (-5 to -10\u00b0C)\nRefrigerant channels freeze PCM\nWater channels chill on pour",
     "Standard component (Kaori, SWEP)\nScales by plate count:\n4 plates = 5L  |  8 plates = 10L"),
    ("Cover / Manifold", "Routes all fluid connections:\n\u2022 Mains water in (top)\n\u2022 Chilled water out (to tap)\n\u2022 Refrigerant lines (to HX)",
     "Only complex manufactured part\nLifts off for full tank access\nInjection-molded + brazed fittings"),
]

for i, (title, desc, note) in enumerate(components):
    x = Inches(0.8) + Inches(4.1) * i

    # Box
    add_shape(slide, x, Inches(1.6), Inches(3.6), Inches(3.8), RGBColor(0x22, 0x22, 0x3A))
    add_shape(slide, x, Inches(1.6), Inches(3.6), Inches(0.06), ACCENT_BLUE)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(1.85), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.word_wrap = False
    set_text(ctf, str(i + 1), size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    tb = add_text_box(slide, x + Inches(0.8), Inches(1.85), Inches(2.6), Inches(0.5))
    set_text(tb.text_frame, title, size=20, color=WHITE, bold=True)

    # Description
    tb = add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3.0), Inches(1.8))
    set_text(tb.text_frame, desc, size=14, color=LIGHT_GRAY)

    # Note box
    add_shape(slide, x + Inches(0.15), Inches(4.2), Inches(3.3), Inches(1.0), RGBColor(0x18, 0x18, 0x28))
    tb = add_text_box(slide, x + Inches(0.3), Inches(4.25), Inches(3.0), Inches(0.9))
    set_text(tb.text_frame, note, size=12, color=ACCENT_CYAN)

# Key callout at bottom
tb = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Key design principle:", size=16, color=WHITE, bold=True)
add_paragraph(tf, "Water moves only when the tap is open (mains pressure driven). No circulation pump. "
              "The cover is the single point of service \u2014 lift it off, and the tank is accessible for cleaning. "
              "The HX comes out with the cover.", size=15, color=LIGHT_GRAY, space_before=Pt(6))

# ============================================================
# SLIDE 5: Cross-Section View
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Inside ChillStream", size=36, color=WHITE, bold=True)

# Add the cross-section image, centered on the slide
import os
img_path = "/Users/yoel/Downloads/unnamed.jpg"
if os.path.exists(img_path):
    # Image is landscape ~1024x600ish. Scale to fill most of the slide area below title.
    img_w = Inches(9.0)
    img_h = Inches(5.5)
    img_left = (W - img_w) // 2
    img_top = Inches(1.4)
    slide.shapes.add_picture(img_path, img_left, img_top, img_w, img_h)

# Subtitle below title
tb = add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4))
set_text(tb.text_frame,
         "Annotated cross-section — all major subsystems visible. Section references map to the system description document.",
         size=13, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6: IoT Platform / Software
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Fleet-First IoT Platform", size=36, color=WHITE, bold=True)

# Left side - fleet dashboard features
add_accent_bar(slide, Inches(0.8), Inches(1.6), height=Inches(4.5), color=ACCENT_CYAN)

tb = add_text_box(slide, Inches(1.2), Inches(1.6), Inches(5.5), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Fleet Dashboard", size=24, color=ACCENT_CYAN, bold=True)

features = [
    ("Real-time monitoring", "Temperatures, capacity %, compressor state across all units"),
    ("Usage analytics", "Pours/day, peak times, demand vs. capacity \u2014 sizing recommendations"),
    ("Predictive maintenance", "Compressor anomaly detection, condenser fouling alerts"),
    ("Filter management", "Fleet-wide filter life tracking, replacement scheduling"),
    ("OTA firmware updates", "Staged rollout, automatic rollback, zero downtime"),
    ("Remote commands", "Reboot, config push, maintenance mode, diagnostics"),
]

for title, desc in features:
    add_paragraph(tf, "", size=6, color=WHITE, space_before=Pt(8))
    add_paragraph(tf, title, size=17, color=WHITE, bold=True, space_before=Pt(2), space_after=Pt(0))
    add_paragraph(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(0))

# Right side - on-unit display
right_x = Inches(7.5)

add_shape(slide, right_x, Inches(1.6), Inches(4.8), Inches(2.2), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, right_x, Inches(1.6), Inches(4.8), Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, right_x + Inches(0.3), Inches(1.8), Inches(4.2), Inches(0.5))
set_text(tb.text_frame, "On-Unit Display", size=20, color=ACCENT_BLUE, bold=True)

tb = add_text_box(slide, right_x + Inches(0.3), Inches(2.3), Inches(4.2), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "\u2022 Cold capacity gauge (% remaining)", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 Live output temperature during pour", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 Filter life indicator", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 Recharge ETA when capacity is low", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 PIN-protected operator settings", size=14, color=LIGHT_GRAY)

# Connectivity box
add_shape(slide, right_x, Inches(4.2), Inches(4.8), Inches(2.2), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, right_x, Inches(4.2), Inches(4.8), Inches(0.06), ORANGE)

tb = add_text_box(slide, right_x + Inches(0.3), Inches(4.4), Inches(4.2), Inches(0.5))
set_text(tb.text_frame, "Connectivity", size=20, color=ORANGE, bold=True)

tb = add_text_box(slide, right_x + Inches(0.3), Inches(4.9), Inches(4.2), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "\u2022 WiFi (ESP32-S3, built-in \u2014 $0 extra HW)", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 MQTT telemetry + HTTPS for OTA/config", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 TLS 1.3, X.509 device certificates", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 Works fully offline \u2014 buffers 7 days locally", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "\u2022 AP mode captive portal for setup", size=14, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6: Market Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Market Opportunity", size=36, color=WHITE, bold=True)

# Market segments
segments = [
    ("Schools (K-12)", "130,000+ in the US alone", "Peak demand at every bell \u2014 the hardest use case.\nBudget-sensitive. Fleet purchases by district.", "$800-1,100"),
    ("Offices", "~1M commercial offices in US", "Steady use with lunch/break clusters.\nFacility managers want zero-maintenance.", "$600-900"),
    ("Healthcare / Clinics", "~230,000 medical offices", "Waiting rooms, staff areas.\nHygiene and compliance matter.", "$800-1,100"),
    ("Hospitality", "Hotels, coworking, airports", "High-visibility locations.\nPremium positioning possible.", "$900-1,200"),
]

for i, (name, size_str, desc, price) in enumerate(segments):
    y = Inches(1.5) + Inches(1.35) * i

    # Row background
    bg_color = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x1E, 0x1E, 0x34)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.2), bg_color)

    # Segment name
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(2.2), Inches(0.5))
    set_text(tb.text_frame, name, size=20, color=ACCENT_CYAN, bold=True)

    # Market size
    tb = add_text_box(slide, Inches(1.0), y + Inches(0.65), Inches(2.2), Inches(0.4))
    set_text(tb.text_frame, size_str, size=13, color=LIGHT_GRAY)

    # Description
    tb = add_text_box(slide, Inches(3.5), y + Inches(0.15), Inches(5.5), Inches(1.0))
    set_text(tb.text_frame, desc, size=14, color=LIGHT_GRAY)

    # Price range
    tb = add_text_box(slide, Inches(9.5), y + Inches(0.25), Inches(2.5), Inches(0.5))
    set_text(tb.text_frame, price, size=22, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

    tb = add_text_box(slide, Inches(9.5), y + Inches(0.7), Inches(2.5), Inches(0.4))
    set_text(tb.text_frame, "per unit", size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# Bottom insight
tb = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.7))
set_text(tb.text_frame,
         "Entry point: Schools (hardest use case = best proof point) and offices (highest volume).",
         size=16, color=ACCENT_CYAN, bold=False)

# ============================================================
# SLIDE 7: Unit Economics & BOM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Unit Economics", size=36, color=WHITE, bold=True)

# BOM breakdown - left side
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5), ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.4))
set_text(tb.text_frame, "5L MODEL \u2014 BOM BREAKDOWN", size=16, color=WHITE, bold=True)

bom_items = [
    ("Cold assembly (tank + HX + cover + PCM)", "$75", "30%", ACCENT_CYAN),
    ("Refrigeration system", "$50", "20%", ACCENT_BLUE),
    ("Hot water system", "$25", "10%", LIGHT_GRAY),
    ("Water path (valves, flow sensor, filter housing)", "$20", "8%", LIGHT_GRAY),
    ("Electronics (ESP32, PCB, sensors, relays, PSU)", "$25", "10%", LIGHT_GRAY),
    ("Display + UI", "$15", "6%", LIGHT_GRAY),
    ("Enclosure + packaging + misc", "$43", "17%", LIGHT_GRAY),
]

for i, (item, cost, pct, color) in enumerate(bom_items):
    y = Inches(2.2) + Inches(0.4) * i

    tb = add_text_box(slide, Inches(1.0), y, Inches(3.8), Inches(0.35))
    set_text(tb.text_frame, item, size=13, color=LIGHT_GRAY)

    tb = add_text_box(slide, Inches(4.9), y, Inches(0.8), Inches(0.35))
    set_text(tb.text_frame, cost, size=14, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

    tb = add_text_box(slide, Inches(5.8), y, Inches(0.6), Inches(0.35))
    set_text(tb.text_frame, pct, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# Total line
y = Inches(2.2) + Inches(0.4) * len(bom_items) + Inches(0.1)
add_shape(slide, Inches(1.0), y, Inches(5.4), Inches(0.01), LIGHT_GRAY)
y += Inches(0.08)

tb = add_text_box(slide, Inches(1.0), y, Inches(3.8), Inches(0.35))
set_text(tb.text_frame, "Materials subtotal", size=14, color=WHITE, bold=True)
tb = add_text_box(slide, Inches(4.9), y, Inches(0.8), Inches(0.35))
set_text(tb.text_frame, "$253", size=14, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

y += Inches(0.35)
tb = add_text_box(slide, Inches(1.0), y, Inches(3.8), Inches(0.35))
set_text(tb.text_frame, "Assembly + test + brazing + scrap (5%)", size=13, color=LIGHT_GRAY)
tb = add_text_box(slide, Inches(4.9), y, Inches(0.8), Inches(0.35))
set_text(tb.text_frame, "$133", size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

y += Inches(0.4)
add_shape(slide, Inches(1.0), y, Inches(5.4), Inches(0.01), ACCENT_CYAN)
y += Inches(0.08)
tb = add_text_box(slide, Inches(1.0), y, Inches(3.8), Inches(0.4))
set_text(tb.text_frame, "TOTAL COGS \u2014 5L UNIT", size=16, color=ACCENT_CYAN, bold=True)
tb = add_text_box(slide, Inches(4.6), y, Inches(1.4), Inches(0.4))
set_text(tb.text_frame, "$386", size=22, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.RIGHT)

# Right side - pricing cards
right_x = Inches(7.3)

# 5L card
add_shape(slide, right_x, Inches(1.5), Inches(2.5), Inches(2.8), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, right_x, Inches(1.5), Inches(2.5), Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, right_x + Inches(0.2), Inches(1.7), Inches(2.1), Inches(0.4))
set_text(tb.text_frame, "5L Model", size=18, color=WHITE, bold=True)

tb = add_text_box(slide, right_x + Inches(0.2), Inches(2.2), Inches(2.1), Inches(0.5))
set_text(tb.text_frame, "$749", size=36, color=ACCENT_BLUE, bold=True)

tb = add_text_box(slide, right_x + Inches(0.2), Inches(2.8), Inches(2.1), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "COGS: $386", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Gross margin: 48%", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Offices, small clinics", size=13, color=LIGHT_GRAY)

# 10L card
add_shape(slide, right_x + Inches(2.8), Inches(1.5), Inches(2.5), Inches(2.8), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, right_x + Inches(2.8), Inches(1.5), Inches(2.5), Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, right_x + Inches(3.0), Inches(1.7), Inches(2.1), Inches(0.4))
set_text(tb.text_frame, "10L Model", size=18, color=WHITE, bold=True)

tb = add_text_box(slide, right_x + Inches(3.0), Inches(2.2), Inches(2.1), Inches(0.5))
set_text(tb.text_frame, "$999", size=36, color=ACCENT_CYAN, bold=True)

tb = add_text_box(slide, right_x + Inches(3.0), Inches(2.8), Inches(2.1), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "COGS: $442", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Gross margin: 56%", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Schools, large offices", size=13, color=LIGHT_GRAY)

# Filter card
add_shape(slide, right_x, Inches(4.6), Inches(5.3), Inches(1.4), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, right_x, Inches(4.6), Inches(5.3), Inches(0.06), ORANGE)

tb = add_text_box(slide, right_x + Inches(0.2), Inches(4.8), Inches(2.5), Inches(0.4))
set_text(tb.text_frame, "Filter Cartridge", size=18, color=ORANGE, bold=True)

tb = add_text_box(slide, right_x + Inches(0.2), Inches(5.2), Inches(2.0), Inches(0.5))
set_text(tb.text_frame, "$20", size=30, color=ORANGE, bold=True)

tb = add_text_box(slide, right_x + Inches(2.8), Inches(4.8), Inches(2.3), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "COGS: $6  |  Margin: 70%", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "2-3 replacements/year per unit", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "NFC tag for auto lifecycle tracking", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Proprietary twist-lock (recurring revenue)", size=13, color=ACCENT_CYAN)

# ============================================================
# SLIDE 8: Recurring Revenue
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ORANGE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Recurring Revenue Model", size=36, color=WHITE, bold=True)

# Revenue streams
add_accent_bar(slide, Inches(0.8), Inches(1.6), height=Inches(3.0), color=ORANGE)

tb = add_text_box(slide, Inches(1.2), Inches(1.6), Inches(5.5), Inches(3.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Two recurring streams per unit:", size=20, color=WHITE, bold=True)

add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Filter cartridges", size=20, color=ORANGE, bold=True, space_before=Pt(12))
add_paragraph(tf, "$20/cartridge \u00d7 2-3/year = $40-60/unit/year", size=16, color=LIGHT_GRAY, space_before=Pt(2))
add_paragraph(tf, "70% gross margin. Proprietary twist-lock with NFC.", size=14, color=LIGHT_GRAY)

add_paragraph(tf, "Cloud SaaS (optional)", size=20, color=ORANGE, bold=True, space_before=Pt(16))
add_paragraph(tf, "$10/month/unit = $120/unit/year", size=16, color=LIGHT_GRAY, space_before=Pt(2))
add_paragraph(tf, "85% gross margin. Fleet dashboard, analytics, OTA.", size=14, color=LIGHT_GRAY)

add_paragraph(tf, "", size=8, color=WHITE)
add_paragraph(tf, "Combined: $160-180/unit/year", size=20, color=WHITE, bold=True, space_before=Pt(12))

# Right side - fleet scaling table
right_x = Inches(7.3)

add_shape(slide, right_x, Inches(1.5), Inches(5.3), Inches(0.5), ACCENT_BLUE)
tb = add_text_box(slide, right_x + Inches(0.2), Inches(1.55), Inches(4.8), Inches(0.4))
set_text(tb.text_frame, "ANNUAL RECURRING AT FLEET SCALE", size=14, color=WHITE, bold=True)

fleet_data = [
    ("10 units", "$1,600-1,800", "Single site (small office building)"),
    ("50 units", "$8,000-9,000", "School district or multi-floor office"),
    ("100 units", "$16,000-18,000", "Regional facility operator"),
    ("500 units", "$80,000-90,000", "National FM company"),
    ("1,000 units", "$160,000-180,000", "Enterprise account"),
]

for i, (fleet, revenue, desc) in enumerate(fleet_data):
    y = Inches(2.2) + Inches(0.65) * i
    bg_color = RGBColor(0x22, 0x22, 0x3A) if i % 2 == 0 else RGBColor(0x1E, 0x1E, 0x34)
    add_shape(slide, right_x, y, Inches(5.3), Inches(0.55), bg_color)

    tb = add_text_box(slide, right_x + Inches(0.2), y + Inches(0.08), Inches(1.2), Inches(0.4))
    set_text(tb.text_frame, fleet, size=15, color=WHITE, bold=True)

    tb = add_text_box(slide, right_x + Inches(1.5), y + Inches(0.08), Inches(1.8), Inches(0.4))
    set_text(tb.text_frame, revenue, size=15, color=ORANGE, bold=True)

    tb = add_text_box(slide, right_x + Inches(3.4), y + Inches(0.08), Inches(1.7), Inches(0.4))
    set_text(tb.text_frame, desc, size=12, color=LIGHT_GRAY)

# Bottom callout
tb = add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(1.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Hardware is the door opener. Recurring revenue is the business model.", size=20, color=WHITE, bold=True)
add_paragraph(tf, "At 1,000 deployed units, annual recurring exceeds initial hardware revenue within 12-18 months of fleet buildout.", size=15, color=LIGHT_GRAY, space_before=Pt(6))

# ============================================================
# SLIDE 9: Competitive Advantage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Why ChillStream Wins", size=36, color=WHITE, bold=True)

advantages = [
    ("Physics advantage",
     "Latent heat stores 5-10x more energy than sensible heat per kg. A 3 kg PCM mass replaces a 15-20L cold water tank. This is why we fit in a countertop form factor.",
     ACCENT_BLUE),
    ("No moving parts in the water path",
     "No pump, no circulator. Water moves on mains pressure only. Fewer failure modes. Lower cost. Silent during dispense.",
     ACCENT_CYAN),
    ("Standard components, novel assembly",
     "The plate HX is off-the-shelf (Kaori, SWEP). The compressor is catalog (Secop, Embraco). The tank is a simple vessel. The innovation is in the assembly and the PCM coating \u2014 not in reinventing the parts.",
     ORANGE),
    ("Fleet-first IoT from day one",
     "Built for operators managing 10-1,000 units. Not a consumer gadget with an app bolted on. Predictive maintenance, sizing recommendations, OTA updates, compliance reporting.",
     ACCENT_BLUE),
    ("Recurring revenue built into the product",
     "Proprietary filter with NFC lifecycle tracking. Cloud SaaS for fleet management. Combined: $160-180/unit/year at 70-85% margin.",
     ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(advantages):
    y = Inches(1.5) + Inches(1.1) * i

    add_accent_bar(slide, Inches(0.8), y + Inches(0.05), height=Inches(0.85), color=color)

    tb = add_text_box(slide, Inches(1.2), y, Inches(3.0), Inches(0.4))
    set_text(tb.text_frame, title, size=18, color=color, bold=True)

    tb = add_text_box(slide, Inches(4.5), y, Inches(8.0), Inches(0.9))
    set_text(tb.text_frame, desc, size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 10: Roadmap & Risks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "Roadmap & Key Risks", size=36, color=WHITE, bold=True)

# Timeline - left side
add_shape(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(0.5), ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.0), Inches(1.55), Inches(6.0), Inches(0.4))
set_text(tb.text_frame, "18-MONTH DEVELOPMENT ROADMAP", size=14, color=WHITE, bold=True)

milestones = [
    ("M1-M3", "PCM + HX prototype", "Material selection, thermal simulation, bench prototype of submerged plate HX"),
    ("M4-M6", "Full functional prototype", "Complete unit: cold + hot + electronics + display + WiFi"),
    ("M7-M9", "EVT", "Engineering validation. Performance testing. Begin certification."),
    ("M10-M12", "DVT", "Design validation. Reliability testing. 10,000-cycle PCM test."),
    ("M13-M14", "Pilot fleet", "5-10 units deployed in real schools and offices. Field data."),
    ("M15-M18", "PVT \u2192 Production", "Manufacturing validation. Certification complete. Ship."),
]

for i, (time, name, desc) in enumerate(milestones):
    y = Inches(2.2) + Inches(0.7) * i

    # Time badge
    badge = add_shape(slide, Inches(0.8), y, Inches(1.0), Inches(0.45),
                      ACCENT_BLUE if i < 2 else ACCENT_CYAN if i < 4 else ORANGE)
    tb_badge = add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.0), Inches(0.35))
    set_text(tb_badge.text_frame, time, size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Milestone name
    tb = add_text_box(slide, Inches(2.0), y, Inches(2.0), Inches(0.35))
    set_text(tb.text_frame, name, size=15, color=WHITE, bold=True)

    # Description
    tb = add_text_box(slide, Inches(4.1), y, Inches(3.2), Inches(0.55))
    set_text(tb.text_frame, desc, size=12, color=LIGHT_GRAY)

# Risks - right side
right_x = Inches(7.8)
add_shape(slide, right_x, Inches(1.5), Inches(4.8), Inches(0.5), RED_ACCENT)
tb = add_text_box(slide, right_x + Inches(0.2), Inches(1.55), Inches(4.4), Inches(0.4))
set_text(tb.text_frame, "TOP RISKS & MITIGATIONS", size=14, color=WHITE, bold=True)

risks = [
    ("PCM encapsulation durability",
     "Can coated plates survive 10,000+ freeze-thaw cycles? \u2192 Material science partnership + accelerated life test by M3"),
    ("HX thermal performance",
     "Can 4 plates achieve 25\u00b0C\u21924\u00b0C at 1 L/min? \u2192 Thermal simulation + bench prototype by M3"),
    ("Cover manifold complexity",
     "Single part routes water + refrigerant + sensors. \u2192 Industrial design + mechanical prototype by M4"),
    ("Mineral fouling",
     "HX plates always submerged in water. \u2192 Upstream filter + descaling protocol + accelerated fouling test"),
]

for i, (risk, mitigation) in enumerate(risks):
    y = Inches(2.2) + Inches(1.2) * i

    add_shape(slide, right_x, y, Inches(4.8), Inches(1.0), RGBColor(0x22, 0x22, 0x3A))

    tb = add_text_box(slide, right_x + Inches(0.2), y + Inches(0.08), Inches(4.4), Inches(0.35))
    set_text(tb.text_frame, risk, size=14, color=RED_ACCENT, bold=True)

    tb = add_text_box(slide, right_x + Inches(0.2), y + Inches(0.42), Inches(4.4), Inches(0.55))
    set_text(tb.text_frame, mitigation, size=12, color=LIGHT_GRAY)

# ============================================================
# SLIDE 11: The Ask
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_CYAN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "The Ask", size=36, color=WHITE, bold=True)

# Phase boxes
phases = [
    ("Phase 1: Prove the Physics", "M1-M6", "$250K",
     [
         "PCM material selection and encapsulation testing",
         "Submerged plate HX thermal prototype",
         "Cover manifold mechanical prototype",
         "Full functional prototype (1 unit)",
         "Thermal performance validation",
     ], ACCENT_BLUE),
    ("Phase 2: Validate the Product", "M7-M14", "$400K",
     [
         "EVT + DVT (10-20 units)",
         "Certification (UL, NSF, FCC/CE)",
         "IoT platform development (dashboard, OTA)",
         "Pilot fleet: 5-10 units in real deployments",
         "Field data collection and iteration",
     ], ACCENT_CYAN),
    ("Phase 3: Scale to Market", "M15-M18+", "$350K",
     [
         "PVT and production tooling",
         "Manufacturing line setup",
         "First 500-unit production run",
         "Sales team and channel partnerships",
         "Launch marketing",
     ], ORANGE),
]

for i, (title, timeline, cost, items, color) in enumerate(phases):
    x = Inches(0.8) + Inches(4.1) * i

    add_shape(slide, x, Inches(1.5), Inches(3.6), Inches(4.8), RGBColor(0x22, 0x22, 0x3A))
    add_shape(slide, x, Inches(1.5), Inches(3.6), Inches(0.06), color)

    # Title
    tb = add_text_box(slide, x + Inches(0.25), Inches(1.7), Inches(3.1), Inches(0.6))
    set_text(tb.text_frame, title, size=16, color=color, bold=True)

    # Timeline + cost
    tb = add_text_box(slide, x + Inches(0.25), Inches(2.3), Inches(1.5), Inches(0.4))
    set_text(tb.text_frame, timeline, size=14, color=LIGHT_GRAY)

    tb = add_text_box(slide, x + Inches(2.0), Inches(2.3), Inches(1.4), Inches(0.4))
    set_text(tb.text_frame, cost, size=22, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

    # Items
    add_shape(slide, x + Inches(0.2), Inches(2.85), Inches(3.2), Inches(0.01), RGBColor(0x33, 0x33, 0x55))

    tb = add_text_box(slide, x + Inches(0.25), Inches(3.0), Inches(3.1), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, items[0], size=13, color=LIGHT_GRAY)
    for item in items[1:]:
        add_paragraph(tf, item, size=13, color=LIGHT_GRAY, space_before=Pt(6))

# Total bar at bottom
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), RGBColor(0x22, 0x22, 0x3A))
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.06), WHITE)

tb = add_text_box(slide, Inches(1.2), Inches(6.6), Inches(5), Inches(0.5))
set_text(tb.text_frame, "Total investment to production:", size=18, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(8.5), Inches(6.55), Inches(3.5), Inches(0.5))
set_text(tb.text_frame, "$1M over 18 months", size=26, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.RIGHT)


# ============================================================
# Save
# ============================================================
output_path = "/Users/yoel/Desktop/product_ideas/chillstream/ChillStream_Executive_Deck.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
